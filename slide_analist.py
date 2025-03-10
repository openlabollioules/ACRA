from pptx import Presentation

def analyze_presentation(file_path="./templates/CRA_LA_GT_TEMPLATE.pptx"):

    # Load the PowerPoint file
    pptx_file = file_path    
    prs = Presentation(pptx_file)
    
    # Print total slides
    print(f"Total Slides: {len(prs.slides)}")

    # Loop through each slide
    for slide_index, slide in enumerate(prs.slides):
        print(f"\n--- Slide {slide_index + 1} ---")

        # List placeholders
        for placeholder in slide.placeholders:
            print(f"  Placeholder {placeholder.placeholder_format.idx}: {placeholder.text}")

        # List all shapes
        for shape_index, shape in enumerate(slide.shapes):
            shape_type = type(shape).__name__
            print(f"  Shape {shape_index} ({shape_type})")

            # If shape has text, print it
            if shape.has_text_frame:
                print(f"    - Text: {shape.text}")

            # If shape is an image
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                print("    - This is an image.")
                
            # If shape is a table
            if shape.has_table:
                print("    - This is a table.")
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    row_text = [cell.text for cell in row.cells]  # Extract text from each cell
                    print(f"    Row {row_idx}: {row_text}")
            
            # If shape is a chart
            if shape.has_chart:
                print(f"  Shape {shape_index} (Chart): {shape.chart.chart_type}")

                # Extract series names and data points
                for series in shape.chart.plots[0].series:
                    series_name = series.name
                    data_points = [pt for pt in series.values]
                    print(f"    Series: {series_name}, Data: {data_points}")
                    

    
