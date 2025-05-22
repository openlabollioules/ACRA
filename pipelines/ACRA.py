import json
import os
import sys
import shutil
import requests
import uuid
import sqlite3
import time
from typing import List, Union, Generator, Iterator, Dict, Any
from langchain_ollama import  OllamaLLM
from dotenv import load_dotenv
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), "..","src")))
from core import summarize_ppt, get_slide_structure, generate_pptx_from_text
from services import merge_pptx, delete_matching_files_in_openwebui

from OLLibrary.utils.text_service import remove_tags_keep, remove_tags_no_keep
from OLLibrary.utils.log_service import setup_logging, get_logger
from OLLibrary.utils.json_service import extract_json

import logging

# Set up the main application logger
setup_logging(app_name="ACRA_Pipeline")
# Use a specific logger for this module
log = get_logger(__name__)
UPLOAD_FOLDER = os.getenv("UPLOAD_FOLDER", "pptx_folder")
OUTPUT_FOLDER = os.getenv("OUTPUT_FOLDER", "OUTPUT")
# Use absolute path for mappings folder
MAPPINGS_FOLDER = os.path.abspath(os.getenv("MAPPINGS_FOLDER", os.path.join(os.getcwd(), "mappings")))

class Pipeline:
    def __init__(self):
        load_dotenv()
        log.info("Initializing ACRA Pipeline")
        self.last_response = None

        # Fix: Properly convert USE_API to boolean
        use_api_env = os.getenv("USE_API", "False")
        self.use_api = use_api_env.lower() in ("true", "1", "t", "yes", "y")
        print(f"USE_API: {self.use_api}")
        # self.model = OllamaLLM(model="deepseek-r1:8b", base_url="http://host.docker.internal:11434", num_ctx=32000)
        self.streaming_model = OllamaLLM(model="qwen3:30b-a3b", base_url="http://host.docker.internal:11434", num_ctx=32000, stream=True)

        self.api_url = "http://host.docker.internal:5050"

        self.openwebui_api = "http://host.docker.internal:3030/api/v1/"
        self.openwebui_db_path = os.getenv("OPENWEBUI_DB_PATH", "./open-webui/webui.db")

        # Réduire le contexte pour le petit modèle et activer le streaming
        self.small_model = OllamaLLM(model="gemma3:latest", base_url="http://host.docker.internal:11434", num_ctx=16000, stream=True)

        self.file_path_list = []
        self.openwebui_api_key = os.getenv("OPENWEBUI_API_KEY")
        if not self.openwebui_api_key:
            log.error("OPENWEBUI_API_KEY is not set")
            raise ValueError("OPENWEBUI_API_KEY is not set")

        self.chat_id = ""
        self.current_chat_id = ""  # To track conversation changes
        self.system_prompt = ""
        self.message_id = 0
        
        # Variable pour stocker la structure traitée
        self.cached_structure = None

        # State tracking
        self.waiting_for_confirmation = False
        self.confirmation_command = ""
        self.confirmation_additional_info = ""
        
        # File ID mapping
        self.file_id_mapping = {}  # Maps {file_path: file_id}
        
        # Create mappings folder if it doesn't exist
        try:
            global MAPPINGS_FOLDER
            log.info(f"Creating mappings directory at {MAPPINGS_FOLDER}")
            os.makedirs(MAPPINGS_FOLDER, exist_ok=True)
            # Test if we can write to the directory
            test_file_path = os.path.join(MAPPINGS_FOLDER, "test_write.txt")
            with open(test_file_path, 'w') as f:
                f.write("Test write access")
            os.remove(test_file_path)
            log.info(f"Mappings directory created and writable at {MAPPINGS_FOLDER}")
        except Exception as e:
            log.error(f"Error creating or accessing mappings directory: {str(e)}")
            # Fallback to using OUTPUT_FOLDER for mappings
            log.warning(f"Using OUTPUT_FOLDER for mappings as fallback")
            MAPPINGS_FOLDER = os.path.abspath(OUTPUT_FOLDER)
        
        log.info("ACRA Pipeline initialized successfully")

    def generate_report(self, foldername, info):
        """
        Génère un rapport à partir du texte fourni en utilisant une requête POST.
        Génère un nouveau fichier avec un timestamp unique à chaque appel.
        
        Args:
            foldername (str): Le nom du dossier où stocker le rapport
            info (str): Le texte à analyser pour générer le rapport
            
        Returns:
            dict: Résultat de la requête avec l'URL de téléchargement
        """
        log.info(f"Generating report for folder: {foldername}")
        log.info(f"use_api setting is: {self.use_api}")
        
        import datetime
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        log.info(f"Creating report with timestamp: {timestamp}")
        
        if self.use_api:
            log.info("Using API endpoint to generate report")
            endpoint = f"generate_report/{foldername}?info={info}&timestamp={timestamp}"
            result = self.fetch(endpoint)
            if "error" in result:
                return result
                
            return self.download_file_openwebui(result["summary"])

        log.info("Using direct function call to generate report")
        result = generate_pptx_from_text(foldername, info, timestamp)
        if "error" in result:
            return result
            
        upload_result = self.download_file_openwebui(result["summary"])
        # Save the mapping after uploading the new file
        self.save_file_mappings()
        return upload_result

    def reset_conversation_state(self):
        """Réinitialise les états spécifiques à une conversation"""
        log.info(f"Resetting conversation state for chat_id: {self.chat_id}")
        self.last_response = None
        self.system_prompt = ""
        self.file_path_list = []
        self.message_id = 0
        self.waiting_for_confirmation = False
        self.confirmation_command = ""
        self.confirmation_additional_info = ""
        self.cached_structure = None
        # Note: We don't clear file_id_mapping here anymore
        # as we want to preserve the mapping between files and OpenWebUI file IDs
        # This mapping is loaded/saved per conversation ID

    def save_file_mappings(self):
        """
        Sauvegarde le mapping des fichiers dans un fichier JSON dans le dossier de mappings.
        """
        try:
            os.makedirs(MAPPINGS_FOLDER, exist_ok=True)
            # Log the mapping file path for debugging
            mapping_file = os.path.join(MAPPINGS_FOLDER, f"{self.chat_id}_file_mappings.json")
            log.info(f"Attempting to save file mappings to {mapping_file}")
            
            # Create directory structure if needed
            os.makedirs(os.path.dirname(mapping_file), exist_ok=True)
            
            # Convert absolute paths to relative for better portability
            relative_mappings = {}
            for file_path, file_id in self.file_id_mapping.items():
                relative_path = os.path.relpath(file_path, os.getcwd())
                relative_mappings[relative_path] = file_id
            
            with open(mapping_file, 'w') as f:
                json.dump(relative_mappings, f)
            
            log.info(f"Successfully saved file mappings to {mapping_file}")
        except Exception as e:
            log.error(f"Error saving file mappings: {str(e)}")
            log.error(f"Current working directory: {os.getcwd()}")
            log.error(f"Mappings directory: {MAPPINGS_FOLDER}")
            log.error(f"Directory exists: {os.path.exists(MAPPINGS_FOLDER)}")
            log.error(f"Directory is writable: {os.access(MAPPINGS_FOLDER, os.W_OK)}")

    def load_file_mappings(self):
        """
        Charge le mapping des fichiers depuis un fichier JSON dans le dossier de mappings.
        """
        try:
            mapping_file = os.path.join(MAPPINGS_FOLDER, f"{self.chat_id}_file_mappings.json")
            log.info(f"Attempting to load file mappings from {mapping_file}")
            
            if os.path.exists(mapping_file):
                with open(mapping_file, 'r') as f:
                    relative_mappings = json.load(f)
                
                # Convert relative paths back to absolute
                self.file_id_mapping = {}
                for relative_path, file_id in relative_mappings.items():
                    abs_path = os.path.abspath(os.path.join(os.getcwd(), relative_path))
                    self.file_id_mapping[abs_path] = file_id
                
                log.info(f"Loaded {len(self.file_id_mapping)} file mappings from {mapping_file}")
            else:
                log.info(f"No mapping file found at {mapping_file}")
                # Try to look in OUTPUT_FOLDER as fallback for legacy mappings
                legacy_mapping_file = os.path.join(OUTPUT_FOLDER, f"{self.chat_id}_file_mappings.json")
                if os.path.exists(legacy_mapping_file):
                    log.info(f"Found legacy mapping file at {legacy_mapping_file}")
                    with open(legacy_mapping_file, 'r') as f:
                        relative_mappings = json.load(f)
                    
                    # Convert relative paths back to absolute
                    self.file_id_mapping = {}
                    for relative_path, file_id in relative_mappings.items():
                        abs_path = os.path.abspath(os.path.join(os.getcwd(), relative_path))
                        self.file_id_mapping[abs_path] = file_id
                    
                    log.info(f"Loaded {len(self.file_id_mapping)} file mappings from legacy location")
                    
                    # Save to new location
                    self.save_file_mappings()
        except Exception as e:
            log.error(f"Error loading file mappings: {str(e)}")
            log.error(f"Current working directory: {os.getcwd()}")
            log.error(f"Mappings directory: {MAPPINGS_FOLDER}")
            log.error(f"Directory exists: {os.path.exists(MAPPINGS_FOLDER)}")
            self.file_id_mapping = {}

    def fetch(self, endpoint):
            """Effectue une requête GET synchrone"""
            url = f"{self.api_url}/{endpoint}"
            log.debug(f"Fetching from: {url}")
            response = requests.get(url)
            if response.status_code != 200:
                log.error(f"API request failed: {response.status_code} - {response.text}")
            return response.json() if response.status_code == 200 else {"error": "Request failed"}

    def post(self, endpoint, data=None, files=None, headers=None):
        """Effectue une requête POST synchrone"""
        # Si l'endpoint commence par http, on le considère comme une URL complète
        print(f"Endpoint: {endpoint}")
        if endpoint.startswith("http"):
            url = endpoint
        else:
            # Sinon on le préfixe avec l'URL de l'API
            url = f"{self.api_url}/{endpoint}"
        log.debug(f"Posting to: {url}")
        response = requests.post(url, data=data, files=files, headers=headers)
        if response.status_code != 200:
            log.error(f"API POST request failed: {response.status_code} - {response.text}")
        return response.json() if response.status_code == 200 else {"error": f"Request failed with status {response.status_code}: {response.text}"}

    # Method to download a filde using the openwebui api. 
    def download_file_openwebui(self, file: str):
        """
        Télécharge un fichier à partir d'un nom de fichier.
        Utilise un système de mapping pour éviter les téléchargements redondants.
        """
        try:
            # Check if we already have this file mapped
            file_path = os.path.abspath(file)
            if file_path in self.file_id_mapping:
                file_id = self.file_id_mapping[file_path]
                log.info(f"File already uploaded, reusing ID: {file_id}")
                download_url = f"http://localhost:3030/api/v1/files/{file_id}/content"
                log.info(f"Reusing existing download URL: {download_url}")
                return {"download_url": download_url}
            
            headers = {
                "accept": "application/json",
                # Remove Content-Type header to let requests set it automatically for multipart/form-data
                "Authorization": f"Bearer {self.openwebui_api_key}"
            }
            url = f"{self.openwebui_api}files/"  # Remove trailing slash
            log.info(f"Uploading file to OpenWebUI API: {url}")
            log.info(f"File path: {file}")
            
            file_id = ""
            try:
                with open(file, "rb") as f:
                    files = {"file": (os.path.basename(file), f, "application/octet-stream")}
                    
                    # Use direct requests instead of self.post for more control
                    response = requests.post(url, headers=headers, files=files)
                    log.info(f"Upload response status: {response.status_code}")
                    
                    if response.status_code != 200:
                        log.error(f"File upload failed: {response.status_code} - {response.text}")
                        return {"error": f"File upload failed: {response.status_code}"}
                    
                    # Parse the response
                    response_data = response.json()
                    log.info(f"Upload response: {response_data}")
                    file_id = response_data.get("id", "")
                    
                    if not file_id:
                        log.error("No file ID returned from upload")
                        return {"error": "No file ID returned from upload"}
                    
                    # Store in our mapping
                    self.file_id_mapping[file_path] = file_id
                    log.info(f"Added file mapping: {file_path} -> {file_id}")
            except Exception as e:
                log.error(f"Error uploading file: {str(e)}")
                return {"error": f"Error uploading file: {str(e)}"}
            
            download_url = f"http://localhost:3030/api/v1/files/{file_id}/content"
            log.info(f"Download URL: {download_url}")
            return {"download_url": download_url}
        except Exception as e:
            log.error(f"Error in download_file_openwebui: {str(e)}")
            return {"error": f"Error in download_file_openwebui: {str(e)}"}
    
    def summarize_folder(self, foldername=None, add_info=None):
        """
        Envoie une demande pour résumer tous les fichiers PowerPoint dans un dossier.
        Génère un nouveau fichier avec un timestamp unique à chaque appel.
        
        Args:
            foldername (str, optional): Le nom du dossier à résumer. Si None, utilise le chat_id.
            add_info (str, optional): Informations supplémentaires à ajouter au résumé.
        Returns:
            dict: Les résultats de l'opération de résumé.
        """
        if foldername is None:
            foldername = self.chat_id
        
        log.info(f"Summarizing folder: {foldername}")
        log.info(f"use_api setting is: {self.use_api}")
        log.info(f"Additional info: {add_info}")
        
        import datetime
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        log.info(f"Creating summary with timestamp: {timestamp}")
        
        if self.use_api:
            log.info("Using API endpoint to summarize folder")
            endpoint = f"acra/{foldername}"
            if add_info:
                endpoint += f"?add_info={add_info}&timestamp={timestamp}"
            else:
                endpoint += f"?timestamp={timestamp}"
            result = self.fetch(endpoint)
            if "error" in result:
                return result
            
            # Let's make sure both the output file and the source files are properly tracked
            # First, track source files from the pptx_folder
            source_folder = os.path.join(UPLOAD_FOLDER, foldername)
            if os.path.exists(source_folder):
                for filename in os.listdir(source_folder):
                    if filename.lower().endswith(".pptx"):
                        source_file_path = os.path.join(source_folder, filename)
                        abs_file_path = os.path.abspath(source_file_path)
                        if abs_file_path not in self.file_id_mapping:
                            log.info(f"Adding source file to tracking: {abs_file_path}")
                            # Will be uploaded and added to mapping when needed
            
            # Then get the download URL for the new summary file
            upload_result = self.download_file_openwebui(result["summary"])
            self.save_file_mappings()  # Save all mappings
            return upload_result
        
        log.info("Using direct function call to summarize folder")
        result = summarize_ppt(foldername, add_info, timestamp)
        if "error" in result:
            return result
        
        # Track both the output file and source files
        upload_result = self.download_file_openwebui(result["summary"])
        
        # Track source files
        source_folder = os.path.join(UPLOAD_FOLDER, foldername)
        if os.path.exists(source_folder):
            for filename in os.listdir(source_folder):
                if filename.lower().endswith(".pptx"):
                    source_file_path = os.path.join(source_folder, filename)
                    abs_file_path = os.path.abspath(source_file_path)
                    if abs_file_path not in self.file_id_mapping:
                        log.info(f"Source file not yet tracked: {abs_file_path}")
                        # Will be uploaded and tracked when needed
        
        # Save the mapping after uploading all files
        self.save_file_mappings()
        return upload_result

    def extract_service_name(self, filename):
        """
        Extrait le nom du service à partir du nom du fichier PowerPoint en utilisant le modèle small_model.
        
        Args:
            filename (str): Le nom du fichier PowerPoint
            
        Returns:
            str: Le nom du service extrait
        """
        prompt = f"Tu es un assistant spécialisé dans le traitement automatique des noms de fichiers. On te donne un nom de fichier de présentation (PowerPoint) contenant un identifiant unique suivi du titre du document. Ton objectif est d'extraire uniquement le titre du document dans un format propre et lisible pour un humain. Le titre est toujours situé après le dernier underscore (`_`) ou après une chaîne d'identifiants. Supprime l'extension `.pptx` ou toute autre extension. Remplace les underscores (`_`) ou tirets (`-`) par des espaces, et capitalise correctement chaque mot. Exemple : **Nom de fichier :** `dc56be63-37a6-4ed6-9223-50f545028ab4_CRA_SERVICE_UX.pptx`   **Titre extrait :** `Service UX` Donne uniquement le titre extrait (pas d'explication), en une seule ligne. voici le nom du fichier : {filename}"
        
        service_name = remove_tags_no_keep(self.small_model.invoke(prompt), '<think>', '</think>')
        # Nettoyer la réponse (enlever les espaces, retours à la ligne, etc.)
        return service_name.strip()

    def analyze_slide_structure(self, foldername=None):
        """
        Analyse la structure des diapositives dans un dossier.
        
        Args:
            foldername (str, optional): Le nom du dossier à analyser. Si None, utilise le chat_id.
        
        Returns:
            dict: Les résultats de l'analyse.
        """
        if foldername is None:
            foldername = self.chat_id
        
        # Ensure foldername is not None after fallback to self.chat_id
        if foldername is None:
            raise Exception("Le nom du dossier et le chat_id sont tous deux None. Impossible d'analyser la structure des diapositives.")
        
        log.info(f"Analyzing slide structure for folder: {foldername}")
        log.info(f"use_api setting is: {self.use_api}")
        
        if self.use_api:
            log.info("Using API endpoint to get slide structure")
            return self.fetch(f"get_slide_structure/{foldername}")
        
        log.info("Using direct function call to get slide structure")
        return get_slide_structure(foldername)
    
    def format_all_slide_data(self, data: dict) -> str:
        """
        Formate les données de plusieurs présentations PPTX en une seule chaîne de texte structurée,
        regroupant tous les projets sans séparation par fichier et avec les événements à venir par service.
        
        Si une structure traitée existe déjà en cache et que data n'est pas None, utilise la structure en cache.
        Sinon, traite la structure et la stocke en cache.

        Args:
            data (dict): Dictionnaire contenant les projets et métadonnées conforme au nouveau format.

        Returns:
            str: Une chaîne de texte structurée listant les informations de tous les projets.
        """
        # Si data est None ou vide, renvoyer un message d'erreur
        if not data:
            return "Aucun fichier PPTX fourni."
            
        # Si data est fourni, mettre à jour le cache
        self.cached_structure = data
        
        # Utiliser la structure en cache si elle existe
        structure_to_process = self.cached_structure
        
        # Vérifier si nous avons des projets
        projects = structure_to_process.get("projects", {})
        if not projects:
            return "Aucun projet trouvé dans les fichiers analysés."
            
        # Récupérer les métadonnées et événements à venir
        metadata = structure_to_process.get("metadata", {})
        processed_files = metadata.get("processed_files", 0)
        upcoming_events = structure_to_process.get("upcoming_events", {})
            
        # Fonction récursive pour afficher les projets à tous les niveaux de hiérarchie
        def format_project_hierarchy(project_name, content, level=0):
            output = ""
            indent = "  " * level
            
            # Format le nom du projet selon son niveau
            if level == 0:
                output += f"{indent}🔶 **{project_name}**\n"
            elif level == 1:
                output += f"{indent}📌 **{project_name}**\n"
            else:
                output += f"{indent}📎 *{project_name}*\n"
            
            # Ajouter les informations si elles existent
            if "information" in content and content["information"]:
                info_lines = content["information"].split('\n')
                for line in info_lines:
                    if line.strip():
                        output += f"{indent}- {line}\n"
                output += "\n"
            
            # Ajouter les alertes critiques
            if "critical" in content and content["critical"]:
                output += f"{indent}- 🔴 **Alertes Critiques:**\n"
                for alert in content["critical"]:
                    output += f"{indent}  - {alert}\n"
                output += "\n"
            
            # Ajouter les alertes à surveiller
            if "small" in content and content["small"]:
                output += f"{indent}- 🟡 **Alertes à surveiller:**\n"
                for alert in content["small"]:
                    output += f"{indent}  - {alert}\n"
                output += "\n"
            
            # Ajouter les avancements
            if "advancements" in content and content["advancements"]:
                output += f"{indent}- 🟢 **Avancements:**\n"
                for advancement in content["advancements"]:
                    output += f"{indent}  - {advancement}\n"
                output += "\n"
            
            # Traiter les sous-projets ou sous-sous-projets de façon récursive
            for key, value in content.items():
                if isinstance(value, dict) and key not in ["information", "critical", "small", "advancements"]:
                    output += format_project_hierarchy(key, value, level + 1)
            
            return output

        # Créer le résultat final
        result = ""
        
        # Afficher le nombre de présentations analysées
        result += f"📊 **Synthèse globale de {processed_files} fichier(s) analysé(s)**\n\n"
        
        # Formater chaque projet principal
        for project_name, project_content in projects.items():
            result += format_project_hierarchy(project_name, project_content)
        
        # Ajouter la section des événements à venir par service
        if upcoming_events:
            result += "\n\n📅 **Événements à venir par service:**\n\n"
            for service, events in upcoming_events.items():
                if events:
                    result += f"- **{service}:**\n"
                    for event in events:
                        result += f"  - {event}\n"
                    result += "\n"
        else:
            result += "\n\n📅 **Événements à venir:** Aucun événement particulier prévu.\n"

        return result.strip()


    def delete_all_files(self, foldername=None):
        """
        Supprime tous les fichiers dans un dossier et met à jour le système de mapping de fichiers.
        Supprime également les fichiers correspondants dans OpenWebUI.
        
        Args:
            foldername (str, optional): Le nom du dossier à vider. Si None, utilise le chat_id.
        
        Returns:
            dict: Les résultats de l'opération de suppression.
        """
        if foldername is None:
            foldername = self.chat_id
        
        log.info(f"Deleting all files for folder: {foldername}")
        
        # Get folder paths for both pptx_folder and OUTPUT
        pptx_folder_path = os.path.join(UPLOAD_FOLDER, foldername)
        output_folder_path = os.path.join(OUTPUT_FOLDER, foldername)
        
        deleted_count = 0
        removed_mappings = 0
        deleted_webui_files = 0
        
        # 1. D'abord, supprimer les fichiers dans OpenWebUI
        log.info(f"Suppression des fichiers OpenWebUI pour le chat: {foldername}")
        webui_result = self.delete_openwebui_files_for_chat(foldername)
        deleted_webui_files = webui_result.get("deleted_count", 0)
        log.info(f"Résultat de la suppression OpenWebUI: {webui_result}")
        
        # Delete files from pptx_folder
        if self.use_api:
            log.info(f"Using API to delete files from {pptx_folder_path}")
            url = f"{self.api_url}/delete_all_pptx_files/{foldername}"
            response = requests.delete(url)
            result = response.json() if response.status_code == 200 else {"error": f"Request failed with status {response.status_code}: {response.text}"}
        else:
            log.info(f"Directly deleting files from {pptx_folder_path}")
            try:
                if os.path.exists(pptx_folder_path):
                    files = os.listdir(pptx_folder_path)
                    for file in files:
                        file_path = os.path.join(pptx_folder_path, file)
                        abs_path = os.path.abspath(file_path)
                        
                        # Remove mapping if exists
                        if abs_path in self.file_id_mapping:
                            del self.file_id_mapping[abs_path]
                            removed_mappings += 1
                            log.info(f"Removed mapping for deleted file: {abs_path}")
                        
                        # Delete the file
                        os.remove(file_path)
                        deleted_count += 1
                        log.info(f"Deleted file: {file_path}")
                        
                    log.info(f"Deleted {deleted_count} files from pptx_folder")
                    result = {"message": f"{deleted_count} fichiers supprimés avec succès."}
                else:
                    result = {"message": "Le dossier n'existe pas."}
            except Exception as e:
                log.error(f"Error deleting files from pptx_folder: {str(e)}")
                result = {"error": f"Erreur lors de la suppression des fichiers: {str(e)}"}
        
        # Also clean up any files in the OUTPUT folder for this conversation
        try:
            if os.path.exists(output_folder_path):
                output_files = os.listdir(output_folder_path)
                output_deleted = 0
                
                for file in output_files:
                    if file.endswith(".pptx"):  # Only delete PPTX files, leave mapping files
                        file_path = os.path.join(output_folder_path, file)
                        abs_path = os.path.abspath(file_path)
                        
                        # Remove mapping if exists
                        if abs_path in self.file_id_mapping:
                            del self.file_id_mapping[abs_path]
                            removed_mappings += 1
                            log.info(f"Removed mapping for deleted output file: {abs_path}")
                        
                        # Delete the file
                        os.remove(file_path)
                        output_deleted += 1
                        log.info(f"Deleted output file: {file_path}")
                
                log.info(f"Deleted {output_deleted} files from OUTPUT folder")
                
                # Update the result message
                if "message" in result:
                    result["message"] += f" Plus {output_deleted} fichiers supprimés du dossier de sortie."
        except Exception as e:
            log.error(f"Error cleaning up OUTPUT folder: {str(e)}")
            # Don't override the main result if there was an error here
        
        # Also clean up any mapping files for this conversation
        try:
            mapping_file = os.path.join(MAPPINGS_FOLDER, f"{foldername}_file_mappings.json")
            if os.path.exists(mapping_file):
                os.remove(mapping_file)
                log.info(f"Deleted mapping file: {mapping_file}")
                result["message"] += f" Fichier de mapping supprimé."
        except Exception as e:
            log.error(f"Error deleting mapping file: {str(e)}")
        
        # Save the updated mappings (empty for this conversation)
        if removed_mappings > 0:
            log.info(f"Removed {removed_mappings} file mappings. Saving updated mappings.")
            self.file_id_mapping = {}  # Clear all mappings for this conversation
            self.save_file_mappings()
        
        # Reset file path list and cached structure
        self.file_path_list = []
        self.cached_structure = None
        
        # Add information about deleted OpenWebUI files
        if "message" in result:
            result["message"] += f" {deleted_webui_files} fichiers supprimés d'OpenWebUI."
        result["deleted_webui_files"] = deleted_webui_files
        
        return result
    
    def get_files_in_folder(self, foldername=None):
        """
        Récupère la liste des fichiers dans un dossier.
        
        Args:
            foldername (str, optional): Le nom du dossier à analyser. Si None, utilise le chat_id.
        
        Returns:
            list: Liste des noms de fichiers PPTX dans le dossier.
        """
        if foldername is None:
            foldername = self.chat_id
            
        folder_path = os.path.join("./pptx_folder", foldername)
        if not os.path.exists(folder_path):
            return []
            
        return [f for f in os.listdir(folder_path) if f.lower().endswith(".pptx")]

    def get_active_conversation_ids(self):
        """
        Récupère tous les IDs de conversation actifs depuis la base de données OpenWebUI.
        
        Returns:
            list: Liste des IDs de conversation actifs.
        """
        conversation_ids = []
        try:
            # Verify database path and log details
            log.info(f"Attempting to access OpenWebUI database at: {self.openwebui_db_path}")
            
            if not os.path.exists(self.openwebui_db_path):
                log.error(f"OpenWebUI database not found at {self.openwebui_db_path}")
                # Try alternative paths
                alt_paths = [
                    "./webui.db",
                    "/app/webui.db",
                    "/app/open-webui/webui.db"
                ]
                for path in alt_paths:
                    if os.path.exists(path):
                        log.info(f"Found database at alternative path: {path}")
                        self.openwebui_db_path = path
                        break
                else:
                    log.error("Could not find OpenWebUI database in any expected location")
                    # Sécurité: Retourner tous les chats existants dans les dossiers pour éviter les suppressions accidentelles
                    return self.get_all_existing_chat_folders()
            
            log.info(f"Connecting to SQLite database at: {self.openwebui_db_path}")
            conn = sqlite3.connect(self.openwebui_db_path)
            cursor = conn.cursor()
            
            # First check if the conversations table exists
            cursor.execute("""
                SELECT name FROM sqlite_master 
                WHERE type='table' AND name='chat'
            """)
            if not cursor.fetchone():
                log.error("Table 'chat' not found in database")
                conn.close()
                # Sécurité: Retourner tous les chats existants dans les dossiers pour éviter les suppressions accidentelles
                return self.get_all_existing_chat_folders()
            
            # Query to get all active conversation IDs
            # Vérifier d'abord si la colonne deleted_at existe
            cursor.execute("PRAGMA table_info(chat)")
            columns = [col[1] for col in cursor.fetchall()]
            
            if "deleted_at" in columns:
                log.info("Using deleted_at column to filter active chats")
                cursor.execute("""
                    SELECT id FROM chat 
                    WHERE deleted_at IS NULL 
                    OR deleted_at = ''
                """)
            else:
                log.info("deleted_at column not found, getting all chat IDs")
                cursor.execute("SELECT id FROM chat")
                
            rows = cursor.fetchall()
            
            base_conversation_ids = [row[0] for row in rows]
            log.info(f"Found {len(base_conversation_ids)} conversations in database: {base_conversation_ids}")
            
            # Vérifier également si les chats ont des messages associés (si la table existe)
            try:
                cursor.execute("""
                    SELECT name FROM sqlite_master 
                    WHERE type='table' AND name='chat'
                """)
                if cursor.fetchone():
                    log.info("Checking messages table for active chats")
                    cursor.execute("""
                        SELECT DISTINCT id FROM chat
                    """)
                    message_chat_ids = [row[0] for row in cursor.fetchall()]
                    log.info(f"Found {len(message_chat_ids)} chats with messages: {message_chat_ids}")
                    
                    # Ajouter les chats qui ont des messages même s'ils ne sont pas dans la liste principale
                    for chat_id in message_chat_ids:
                        if chat_id not in base_conversation_ids:
                            base_conversation_ids.append(chat_id)
                            log.info(f"Added chat with messages that was not in main list: {chat_id}")
            except Exception as e:
                log.error(f"Error checking messages table: {str(e)}")
            
            conn.close()
            
            # Convertir les UUID en strings si nécessaire
            conversation_ids = [str(id) for id in base_conversation_ids]
            log.info(f"Final active conversations: {conversation_ids}")
        except Exception as e:
            log.error(f"Error retrieving conversation IDs from database: {str(e)}")
            log.exception("Database access exception details:")
            # Sécurité: Retourner tous les chats existants dans les dossiers pour éviter les suppressions accidentelles
            return self.get_all_existing_chat_folders()
        
        # Si nous n'avons trouvé aucun chat actif, considérer tous les dossiers existants comme actifs par sécurité
        if not conversation_ids:
            log.warning("No active conversations found in database! Using existing folders as fallback.")
            conversation_ids = self.get_all_existing_chat_folders()
        
        # Toujours inclure le chat actuel
        if self.chat_id and self.chat_id not in conversation_ids:
            log.warning(f"Current chat_id {self.chat_id} not found in active list! Adding it.")
            conversation_ids.append(self.chat_id)
        
        return conversation_ids
    
    def get_all_existing_chat_folders(self):
        """
        Récupère les IDs de tous les dossiers de chat existants dans UPLOAD_FOLDER et OUTPUT_FOLDER.
        Utilisé comme fallback de sécurité en cas d'échec de la détection des chats actifs.
        
        Returns:
            list: Liste des IDs de dossiers de chat existants
        """
        folder_ids = set()
        
        # Récupérer tous les dossiers dans UPLOAD_FOLDER
        if os.path.exists(UPLOAD_FOLDER):
            try:
                for folder_name in os.listdir(UPLOAD_FOLDER):
                    if os.path.isdir(os.path.join(UPLOAD_FOLDER, folder_name)):
                        folder_ids.add(folder_name)
            except Exception as e:
                log.error(f"Error listing UPLOAD_FOLDER: {str(e)}")
        
        # Récupérer tous les dossiers dans OUTPUT_FOLDER
        if os.path.exists(OUTPUT_FOLDER):
            try:
                for folder_name in os.listdir(OUTPUT_FOLDER):
                    if os.path.isdir(os.path.join(OUTPUT_FOLDER, folder_name)):
                        folder_ids.add(folder_name)
            except Exception as e:
                log.error(f"Error listing OUTPUT_FOLDER: {str(e)}")
        
        # Récupérer tous les fichiers de mapping
        if os.path.exists(MAPPINGS_FOLDER):
            try:
                for filename in os.listdir(MAPPINGS_FOLDER):
                    if filename.endswith("_file_mappings.json"):
                        chat_id = filename.split("_file_mappings.json")[0]
                        folder_ids.add(chat_id)
            except Exception as e:
                log.error(f"Error listing MAPPINGS_FOLDER: {str(e)}")
        
        folder_list = list(folder_ids)
        log.info(f"Found {len(folder_list)} existing chat folders as safety fallback: {folder_list}")
        return folder_list

    def cleanup_orphaned_conversations(self):
        """
        Nettoie les dossiers et fichiers de conversations qui n'existent plus dans OpenWebUI.
        Appelé quand le chat_id change pour s'assurer que les ressources sont bien gérées.
        
        Returns:
            dict: Résultats de l'opération de nettoyage
        """
        log.info("Starting cleanup of orphaned conversations")
        active_conversations = self.get_active_conversation_ids()
        
        # Log all active conversations for debugging
        log.info(f"Active conversations: {active_conversations}")
        
        if not active_conversations:
            log.warning("No active conversations found or unable to get conversation list")
            return {"status": "warning", "message": "Could not retrieve active conversations - No cleanup performed for safety", "action": "none"}
        
        # Mesure de sécurité supplémentaire - Si on a moins de 2 conversations actives, on ne nettoie rien
        # car cela pourrait indiquer une erreur dans la détection des chats actifs
        if len(active_conversations) < 2:
            log.warning(f"Only {len(active_conversations)} active conversation(s) found. Skipping cleanup for safety.")
            return {
                "status": "warning", 
                "message": f"Only {len(active_conversations)} active conversation(s) found. Skipping cleanup for safety.", 
                "action": "none"
            }
        
        # Make sure current chat_id is considered active
        if self.chat_id and self.chat_id not in active_conversations:
            log.info(f"Adding current chat_id to active list: {self.chat_id}")
            active_conversations.append(self.chat_id)
        
        deleted_folders = 0
        deleted_files = 0
        deleted_mappings = 0
        
        # Get list of all folders in UPLOAD_FOLDER and OUTPUT_FOLDER
        upload_folders = []
        output_folders = []
        
        try:
            if os.path.exists(UPLOAD_FOLDER):
                upload_folders = [d for d in os.listdir(UPLOAD_FOLDER) 
                                if os.path.isdir(os.path.join(UPLOAD_FOLDER, d))]
                log.info(f"Found {len(upload_folders)} folders in UPLOAD_FOLDER: {upload_folders}")
            
            if os.path.exists(OUTPUT_FOLDER):
                output_folders = [d for d in os.listdir(OUTPUT_FOLDER) 
                               if os.path.isdir(os.path.join(OUTPUT_FOLDER, d))]
                log.info(f"Found {len(output_folders)} folders in OUTPUT_FOLDER: {output_folders}")
        except Exception as e:
            log.error(f"Error listing folders: {str(e)}")
        
        # Vérifier si tous les dossiers sont considérés comme actifs
        # Si c'est le cas, cela pourrait indiquer un problème avec la détection des chats actifs
        folders_to_delete = [f for f in upload_folders if f not in active_conversations]
        if len(folders_to_delete) == 0 and len(upload_folders) > 0:
            log.warning("All upload folders are considered active - this might indicate an issue. Skipping cleanup for safety.")
            return {
                "status": "warning", 
                "message": "All folders are considered active. Skipping cleanup for safety.", 
                "active_conversations": active_conversations,
                "upload_folders": upload_folders,
                "action": "none"
            }
        
        # Clean pptx_folder
        try:
            if os.path.exists(UPLOAD_FOLDER):
                log.info(f"Checking for orphaned folders in {UPLOAD_FOLDER}")
                for folder_name in upload_folders:
                    if folder_name not in active_conversations:
                        folder_path = os.path.join(UPLOAD_FOLDER, folder_name)
                        log.info(f"Deleting orphaned PPTX folder: {folder_path}")
                        try:
                            # Count files before deletion
                            file_count = len([f for f in os.listdir(folder_path) 
                                          if os.path.isfile(os.path.join(folder_path, f))])
                            deleted_files += file_count
                            log.info(f"Deleting {file_count} files from {folder_path}")
                            
                            # Try to forcefully remove the directory
                            shutil.rmtree(folder_path, ignore_errors=True)
                            deleted_folders += 1
                            log.info(f"Successfully deleted folder: {folder_path}")
                        except Exception as e:
                            log.error(f"Failed to delete folder {folder_path}: {str(e)}")
        except Exception as e:
            log.error(f"Error cleaning up pptx folders: {str(e)}")
        
        # Clean OUTPUT folder
        try:
            if os.path.exists(OUTPUT_FOLDER):
                log.info(f"Checking for orphaned folders in {OUTPUT_FOLDER}")
                for folder_name in output_folders:
                    if folder_name not in active_conversations:
                        folder_path = os.path.join(OUTPUT_FOLDER, folder_name)
                        log.info(f"Deleting orphaned OUTPUT folder: {folder_path}")
                        try:
                            # Count files before deletion
                            file_count = len([f for f in os.listdir(folder_path) 
                                          if os.path.isfile(os.path.join(folder_path, f))])
                            deleted_files += file_count
                            log.info(f"Deleting {file_count} files from {folder_path}")
                            
                            # Try to forcefully remove the directory
                            shutil.rmtree(folder_path, ignore_errors=True)
                            deleted_folders += 1
                            log.info(f"Successfully deleted folder: {folder_path}")
                        except Exception as e:
                            log.error(f"Failed to delete folder {folder_path}: {str(e)}")
        except Exception as e:
            log.error(f"Error cleaning up output folders: {str(e)}")
        
        # Clean mappings folder
        try:
            if os.path.exists(MAPPINGS_FOLDER):
                log.info(f"Checking for orphaned mapping files in {MAPPINGS_FOLDER}")
                for filename in os.listdir(MAPPINGS_FOLDER):
                    if filename.endswith("_file_mappings.json"):
                        # Extract chat_id from filename (chat_id_file_mappings.json)
                        chat_id = filename.split("_file_mappings.json")[0]
                        if chat_id not in active_conversations:
                            mapping_path = os.path.join(MAPPINGS_FOLDER, filename)
                            log.info(f"Deleting orphaned mapping file: {mapping_path}")
                            try:
                                os.remove(mapping_path)
                                deleted_mappings += 1
                                log.info(f"Successfully deleted mapping file: {mapping_path}")
                            except Exception as e:
                                log.error(f"Failed to delete mapping file {mapping_path}: {str(e)}")
        except Exception as e:
            log.error(f"Error cleaning up mapping files: {str(e)}")
        
        result = {
            "status": "success",
            "deleted_folders": deleted_folders,
            "deleted_files": deleted_files,
            "deleted_mappings": deleted_mappings,
            "active_conversations": len(active_conversations),
            "upload_folders": upload_folders,
            "output_folders": output_folders,
            "action": "cleanup"
        }
        
        log.info(f"Cleanup results: {result}")
        return result

    async def inlet(self, body: dict, user: dict) -> dict:
        log.info(f"Received body: {body}")
        log.info(f"Metadata: {body.get('metadata', {})}")
        
        # Debug log the current state
        log.info(f"Current state - self.chat_id: '{self.chat_id}', self.current_chat_id: '{self.current_chat_id}'")
        
        # Get conversation ID from body
        new_chat_id = None
        if "metadata" in body and "chat_id" in body["metadata"]:
            new_chat_id = body["metadata"]["chat_id"]
            log.info(f"New chat_id extracted from request metadata: '{new_chat_id}'")
            
            # Check if this is actually a change
            is_change = new_chat_id != self.chat_id
            has_previous = bool(self.chat_id)
            
            log.info(f"Chat ID change detection: is_change={is_change}, has_previous={has_previous}")
            
            # If chat_id changed, we need to save current mappings and load new ones
            if is_change and has_previous:
                log.info(f"*** CHAT ID CHANGED *** from '{self.chat_id}' to '{new_chat_id}'")
                
                # Save current mappings before changing IDs
                self.save_file_mappings()  # Save mappings for old chat
                old_chat_id = self.chat_id
                
                # Update chat_id first to ensure cleanup knows the current one is active
                self.chat_id = new_chat_id
                
                # Get active chats for cleanup
                active_chats = self.get_active_conversation_ids()
                log.info(f"Active chats for cleanup: {active_chats}")
                
                # Directly delete old folders if they're not in active chats
                if old_chat_id not in active_chats and old_chat_id != new_chat_id:
                    log.info(f"Directly cleaning up old chat ID folders: {old_chat_id}")
                    
                    # Supprimer d'abord les fichiers OpenWebUI
                    log.info(f"Suppression des fichiers OpenWebUI pour le chat: {old_chat_id}")
                    webui_result = self.delete_openwebui_files_for_chat(old_chat_id)
                    log.info(f"Résultat de la suppression OpenWebUI: {webui_result}")
                    
                    # Delete UPLOAD_FOLDER for old chat
                    old_upload_folder = os.path.join(UPLOAD_FOLDER, old_chat_id)
                    if os.path.exists(old_upload_folder) and os.path.isdir(old_upload_folder):
                        log.info(f"Deleting old upload folder: {old_upload_folder}")
                        try:
                            shutil.rmtree(old_upload_folder)
                            log.info(f"Successfully deleted: {old_upload_folder}")
                        except Exception as e:
                            log.error(f"Error deleting {old_upload_folder}: {str(e)}")
                    
                    # Delete OUTPUT_FOLDER for old chat
                    old_output_folder = os.path.join(OUTPUT_FOLDER, old_chat_id)
                    if os.path.exists(old_output_folder) and os.path.isdir(old_output_folder):
                        log.info(f"Deleting old output folder: {old_output_folder}")
                        try:
                            shutil.rmtree(old_output_folder)
                            log.info(f"Successfully deleted: {old_output_folder}")
                        except Exception as e:
                            log.error(f"Error deleting {old_output_folder}: {str(e)}")
                    
                    # Delete mapping file for old chat
                    old_mapping_file = os.path.join(MAPPINGS_FOLDER, f"{old_chat_id}_file_mappings.json")
                    if os.path.exists(old_mapping_file):
                        log.info(f"Deleting old mapping file: {old_mapping_file}")
                        try:
                            os.remove(old_mapping_file)
                            log.info(f"Successfully deleted: {old_mapping_file}")
                        except Exception as e:
                            log.error(f"Error deleting {old_mapping_file}: {str(e)}")
                
                # Run the general cleanup as well
                log.info(f"Running general cleanup for chat ID change")
                try:
                    cleanup_result = self.cleanup_orphaned_conversations()
                    log.info(f"Cleanup results: {cleanup_result}")
                    
                    # Force cleanup as a fallback with only current chat ID preserved
                    log.info(f"Performing forced cleanup to ensure old folders are removed")
                    force_result = self.force_cleanup_old_folders([new_chat_id])
                    log.info(f"Force cleanup results: {force_result}")
                except Exception as e:
                    log.error(f"Cleanup process failed: {str(e)}")
                    log.exception("Cleanup exception details:")
                
                # Reset state but preserve file mappings
                self.reset_conversation_state()
                
                # Load mappings for new conversation
                self.load_file_mappings()
            elif not self.chat_id and new_chat_id:
                # First time setting chat_id
                log.info(f"Setting initial chat_id to {new_chat_id}")
                self.chat_id = new_chat_id
                self.load_file_mappings()  # Try to load any existing mappings
                
            if not self.current_chat_id and new_chat_id:
                self.current_chat_id = new_chat_id
        else:
            log.warning("No chat_id found in metadata!")

        # Create folders for the conversation
        # Create pptx_folder for source files
        conversation_folder = os.path.join(UPLOAD_FOLDER, self.chat_id)
        os.makedirs(conversation_folder, exist_ok=True)
        log.info(f"Created/verified upload folder at: {conversation_folder}")
        
        # Create output folder for generated files
        output_folder = os.path.join(OUTPUT_FOLDER, self.chat_id) 
        os.makedirs(output_folder, exist_ok=True)
        log.info(f"Created/verified output folder at: {output_folder}")
        
        # Create mappings folder if needed
        os.makedirs(MAPPINGS_FOLDER, exist_ok=True)

        # Extract files from body['metadata']['files']
        files = body.get("metadata", {}).get("files", [])
        if files:
            # Réinitialiser la structure en cache car de nouveaux fichiers ont été ajoutés
            self.cached_structure = None
            
            # Traiter les fichiers
            for file_entry in files:
                file_data = file_entry.get("file", {})
                filename = file_data.get("filename", "N/A")
                file_id = file_data.get("id", "N/A")

                filecomplete_name = file_id + "_" + filename

                source_path = os.path.join("./uploads", filecomplete_name)
                # Update destination to use conversation foldername
                destination_path = os.path.join(conversation_folder, filecomplete_name)
                
                self.file_path_list.append(destination_path)
                shutil.copy(source_path, destination_path)
                
                # Add to file mappings if not already present
                abs_path = os.path.abspath(destination_path)
                if abs_path not in self.file_id_mapping:
                    self.file_id_mapping[abs_path] = file_id
                    log.info(f"Added mapping for uploaded file: {abs_path} -> {file_id}")
                
                # Extraire et afficher le nom du service pour information
                service_name = self.extract_service_name(filename)
                log.info(f"File {filename} identified as service: {service_name}")
                
            # Analyser la structure
            response = self.analyze_slide_structure(self.chat_id)
            if "error" in response:
                response = f"Erreur lors de l'analyse de la structure: {response['error']}"
            else:
                # Formater la réponse
                response = self.format_all_slide_data(response)
                # Stocker la structure en cache
                self.cached_structure = response
                
            self.system_prompt = "# Voici les informations des fichiers PPTX toutes les informations sont importantes pour la compréhension du message de l'utilisateur et les données sont triées : \n\n" +  response + "# voici le message de l'utilisateur : " 
            
            # Save file mappings after processing new files
            self.save_file_mappings()
        
        return body

    def get_existing_summaries(self, folder_name=None):
        """
        Récupère la liste des fichiers de résumé existants pour le chat_id actuel et les télécharge vers OpenWebUI.
        Cherche les fichiers dans les dossiers OUTPUT et pptx_folder.
        
        Args:
            folder_name (str, optional): Le nom du dossier à analyser. Si None, utilise le chat_id.
        
        Returns:
            list: Liste des tuples (filename, url) des résumés.
        """
        if folder_name is None:
            folder_name = self.chat_id
        log.info(f"Getting existing summaries for folder: {folder_name}")
        
        summaries = []
        
        # Check in the OUTPUT folder first (traditional location for summaries)
        output_path = os.path.join(OUTPUT_FOLDER, folder_name)
        log.info(f"Checking output path: {output_path}")
        
        # Ensure OUTPUT directory exists
        os.makedirs(output_path, exist_ok=True)
        
        # Also check in pptx_folder (new location for reports)
        pptx_path = os.path.join(UPLOAD_FOLDER, folder_name)
        log.info(f"Checking pptx folder path: {pptx_path}")
        
        # Ensure pptx directory exists
        os.makedirs(pptx_path, exist_ok=True)
        
        try:
            # Check for PPTX files in OUTPUT directory
            if os.path.exists(output_path):
                files = os.listdir(output_path)
                log.info(f"Files in output folder: {files}")
                
                for filename in files:
                    if filename.endswith(".pptx"):
                        log.info(f"Found summary file in OUTPUT: {filename}")
                        file_path = os.path.join(output_path, filename)
                        
                        # Upload file to OpenWebUI to get a download URL
                        upload_result = self.download_file_openwebui(file_path)
                        
                        if "error" in upload_result:
                            log.error(f"Error uploading file to OpenWebUI: {upload_result['error']}")
                            continue
                        
                        download_url = upload_result.get("download_url", "")
                        if download_url:
                            log.info(f"Generated OpenWebUI URL: {download_url}")
                            summaries.append(("OUTPUT/" + filename, download_url))
                        else:
                            log.error(f"No download URL received for file: {filename}")
            
            # Check for PPTX files in pptx_folder directory
            if os.path.exists(pptx_path):
                files = os.listdir(pptx_path)
                log.info(f"Files in pptx folder: {files}")
                
                for filename in files:
                    if filename.endswith(".pptx") and ("_summary_" in filename or "_text_summary_" in filename):
                        log.info(f"Found summary/report file in pptx_folder: {filename}")
                        file_path = os.path.join(pptx_path, filename)
                        
                        # Upload file to OpenWebUI to get a download URL
                        upload_result = self.download_file_openwebui(file_path)
                        
                        if "error" in upload_result:
                            log.error(f"Error uploading file to OpenWebUI: {upload_result['error']}")
                            continue
                        
                        download_url = upload_result.get("download_url", "")
                        if download_url:
                            log.info(f"Generated OpenWebUI URL: {download_url}")
                            summaries.append(("pptx_folder/" + filename, download_url))
                        else:
                            log.error(f"No download URL received for file: {filename}")
            
            # Save any new mappings we've created
            if len(summaries) > 0:
                self.save_file_mappings()
                
            log.info(f"Final summaries list: {summaries}")
        except Exception as e:
            log.error(f"Error listing summary files: {str(e)}")
            log.error(f"Current working directory: {os.getcwd()}")
        
        return summaries

    def pipe(self, body: dict, user_message: str, model_id: str, messages: List[dict]) -> Generator[str, None, None]:
        """
        Gère le pipeline de traitement des messages et des commandes spécifiques.

        Cette méthode traite différentes commandes comme /summarize, /structure, et /clear, 
        et gère le streaming de réponses du modèle.

        Args:
            body (dict): Le corps de la requête contenant des métadonnées.
            user_message (str): Le message de l'utilisateur.
            model_id (str): L'identifiant du modèle utilisé.
            messages (List[dict]): Liste des messages précédents.

        Yields:
            str: Réponses formatées en Server-Sent Events (SSE) compatibles avec OpenWebUI.

        Commandes supportées:
        - /summarize: Tente de résumer les fichiers PPTX
        - /structure: Analyse la structure des diapositives
        - /clear: Supprime tous les fichiers de la conversation
        """
        message = user_message.lower()  # Convertir en minuscules pour simplifier la correspondance
        __event_emitter__ = body.get("__event_emitter__")

        # Check if we're waiting for confirmation
        if self.waiting_for_confirmation:
            if message in ["yes", "y", "oui", "o"]:
                self.waiting_for_confirmation = False
                
                # If we were waiting for summarize confirmation
                if self.confirmation_command == "summarize":
                    # Generate a new summary
                    log.info(f"Generating summary with additional info: {self.confirmation_additional_info}")
                    response = self.summarize_folder(add_info=self.confirmation_additional_info)
                    if "error" in response:
                        response = f"Erreur lors de la génération du résumé: {response['error']}"
                    else:
                        response = f"Le résumé de tous les fichiers a été généré avec succès. URL de téléchargement: \n{response.get('download_url', 'Non disponible')}"
                    
                    # Save mappings after generating a new summary
                    self.save_file_mappings()
                    
                    yield f"data: {json.dumps({'choices': [{'message': {'content': response}}]})}\n\n"
                    yield f"data: {json.dumps({'choices': [{'finish_reason': 'stop'}]})}\n\n"
                    self.last_response = response
                    return
            
            elif message in ["no", "n", "non"]:
                self.waiting_for_confirmation = False
                response = "Génération de résumé annulée."
                yield f"data: {json.dumps({'choices': [{'message': {'content': response}}]})}\n\n"
                yield f"data: {json.dumps({'choices': [{'finish_reason': 'stop'}]})}\n\n"
                self.last_response = response
                return
            
            # Reset if we get any other input
            self.waiting_for_confirmation = False
        
        # Gestion des commandes spécifiques
        if "/summarize" in message:
            # Extract additional information after the /summarize command
            additional_info = None
            if " " in message:
                command_parts = message.split(" ", 1)
                if len(command_parts) > 1 and command_parts[1].strip():
                    additional_info = command_parts[1].strip()
            
            # Get existing summaries
            existing_summaries = self.get_existing_summaries()
            log.info(f"ACRA - Pipeline: Existing summaries: {existing_summaries}")
            
            if existing_summaries:
                response = "Voici les résumés existants pour cette conversation:\n\n"
                for filename, url in existing_summaries:
                    response += f"- {filename}: {url}\n"
                
                response += "\nVoulez-vous générer un nouveau résumé? (Oui/Non)"
                
                # Set state to wait for confirmation
                self.waiting_for_confirmation = True
                self.confirmation_command = "summarize"
                self.confirmation_additional_info = additional_info
                
                yield f"data: {json.dumps({'choices': [{'message': {'content': response}}]})}\n\n"
                yield f"data: {json.dumps({'choices': [{'finish_reason': 'stop'}]})}\n\n"
                self.last_response = response
                return
            else:
                # No existing summaries, generate one directly
                response = self.summarize_folder(add_info=additional_info)
                if "error" in response:
                    response = {"error": f"Erreur lors de la génération du résumé: {response['error']}"}
                else:
                    introduction_prompt = f"""Tu es un assistant qui va générer une introduction pour un enssemble de fichiers PPTX je veux juste une description globale des fichiers impliqués dans le message de 
                l'utilisateur pas de cas par cas et sourtout quelque chose de consit et renvoie uniquement l'introduction (pas d'explication) si tu vois une information importante ou une alerte critique, tu dois 
                la signaler dans l'introduction. Voici le contenu de tous les fichiers : {self.system_prompt} Tu dois renvoyer uniquement l'introduction (pas d'explication).
                """
                    introduction = self.small_model.invoke(introduction_prompt if "introduction_prompt" in locals() else self.system_prompt)
                    response = f"{introduction}\n\n Le résumé de tous les fichiers a été généré avec succès.\n\n  ### URL de téléchargement: \n{response.get('download_url', 'Non disponible')}"
                    
                    # Save mappings after generating a new summary
                    self.save_file_mappings()
                
                yield f"data: {json.dumps({'choices': [{'message': {'content': response}}]})}\n\n"
                yield f"data: {json.dumps({'choices': [{'finish_reason': 'stop'}]})}\n\n"
                self.last_response = response
                return
        
        elif "/structure" in message:
            if self.cached_structure is None:
                # Récupérer la structure des diapositives
                response = self.analyze_slide_structure(self.chat_id)
                
                if "error" in response:
                    response_text = f"Erreur lors de l'analyse de la structure: {response['error']}"
                    if __event_emitter__:
                        __event_emitter__({"type": "content", "content": response_text})
                    yield f"data: {json.dumps({'choices': [{'message': {'content': response_text}}]})}\n\n"
                    yield f"data: {json.dumps({'choices': [{'finish_reason': 'stop'}]})}\n\n"
                    self.last_response = response_text
                    return
                
                # Formater les données de la structure
                formatted_response = self.format_all_slide_data(response)
                self.cached_structure = formatted_response
                if __event_emitter__:
                    __event_emitter__({"type": "content", "content": formatted_response})
                yield f"data: {json.dumps({'choices': [{'message': {'content': formatted_response}}]})}\n\n"
                yield f"data: {json.dumps({'choices': [{'finish_reason': 'stop'}]})}\n\n"
                self.last_response = formatted_response
                return
            else:
                if __event_emitter__:
                    __event_emitter__({"type": "content", "content": self.cached_structure})
                yield f"data: {json.dumps({'choices': [{'message': {'content': self.cached_structure}}]})}\n\n"
                yield f"data: {json.dumps({'choices': [{'finish_reason': 'stop'}]})}\n\n"
                self.last_response = self.cached_structure
                return
        elif "/generate" in message:
            # Extraire le texte après la commande
            text_content = user_message.replace("/generate", "").strip()
            if not text_content:
                response = "Veuillez fournir du texte après la commande /generate pour générer un rapport."
            else:
                # On utilise la méthode generate_report qui maintenant fait un POST avec le texte dans le body
                response = self.generate_report(self.chat_id, text_content)
                if "error" in response:
                    response = f"Erreur lors de la génération du rapport: {response['error']}"
                else:
                    response = f"Le rapport a été généré avec succès à partir du texte fourni.\n\n### URL de téléchargement:\n{response.get('download_url', 'Non disponible')}"
                    # Save mappings after generating a new report
                    self.save_file_mappings()
            
            if __event_emitter__:
                __event_emitter__({"type": "content", "content": response})
            yield f"data: {json.dumps({'choices': [{'message': {'content': response}}]})}\n\n"
            yield f"data: {json.dumps({'choices': [{'finish_reason': 'stop'}]})}\n\n"
            self.last_response = response
            return
        
        elif "/clear" in message:
            log.info("Processing /clear command (using force cleanup method)")
            try:
                # Preserve current chat ID
                preserve_ids = [self.chat_id] if self.chat_id else []
                
                # Extract additional IDs to preserve if specified
                if " " in message:
                    additional_ids = message.split(" ", 1)[1].strip().split()
                    if additional_ids:
                        log.info(f"Additional IDs to preserve: {additional_ids}")
                        preserve_ids.extend(additional_ids)
                
                # Run the standard cleanup
                log.info("Running standard cleanup")
                cleanup_result = self.cleanup_orphaned_conversations()
                
                # Run the forced cleanup, préservant le chat actuel et les IDs supplémentaires
                log.info(f"Running forced cleanup via /clear command, preserving: {preserve_ids}")
                force_result = self.force_cleanup_old_folders(preserve_ids)
                
                # Reset important state variables
                self.file_path_list = []
                self.cached_structure = None
                self.file_id_mapping = {}  # Clear all mappings for this conversation
                
                response = f"Nettoyage terminé!\n\nProtégés: {preserve_ids}\nSupprimés: {force_result['deleted_folders']} dossiers, {force_result['deleted_files']} fichiers locaux, {force_result['deleted_webui_files']} fichiers OpenWebUI"
                log.info(f"Clear command completed successfully. Response: {response}")
            except Exception as e:
                log.error(f"Error processing /clear command: {str(e)}")
                log.exception("Clear command exception details:")
                response = f"Une erreur s'est produite lors du nettoyage: {str(e)}"
                
            if __event_emitter__:
                __event_emitter__({"type": "content", "content": response})
            yield f"data: {json.dumps({'choices': [{'message': {'content': response}}]})}\n\n"
            yield f"data: {json.dumps({'choices': [{'finish_reason': 'stop'}]})}\n\n"
            self.last_response = response
            return

        elif "/merge" in message:
            output_merge = "./OUTPUT/"+self.chat_id + "/merged/" 
            input_merge = "./pptx_folder/" + self.chat_id
            merge_result = merge_pptx(input_merge, output_merge)
            
            if "error" in merge_result:
                response = f"Erreur lors de la fusion des fichiers: {merge_result['error']}"
            else:
                # Get the merged file path from the result
                merged_file = merge_result.get("merged_file")
                if merged_file:
                    # Upload the merged file to OpenWebUI and get download URL
                    upload_result = self.download_file_openwebui(merged_file)
                    if "error" in upload_result:
                        response = f"Les fichiers ont été fusionnés avec succès, mais une erreur s'est produite lors de la génération du lien de téléchargement: {upload_result['error']}"
                    else:
                        response = f"Les fichiers ont été fusionnés avec succès.\n\n### URL de téléchargement:\n{upload_result.get('download_url', 'Non disponible')}"
                else:
                    response = "Les fichiers ont été fusionnés avec succès, mais le chemin du fichier fusionné n'a pas été trouvé."
                
                # Save mappings after uploading the merged file
                self.save_file_mappings()
            
            yield f"data: {json.dumps({'choices': [{'message': {'content': response}}]})}\n\n"
            yield f"data: {json.dumps({'choices': [{'finish_reason': 'stop'}]})}\n\n"
            self.last_response = response
            return

        elif "/regroup" in message:
            # First, get the structure of all files
            if self.cached_structure is None:
                structure_result = self.analyze_slide_structure(self.chat_id)
                if "error" in structure_result:
                    response = f"Erreur lors de l'analyse de la structure: {structure_result['error']}"
            else:
                # Optimiser le prompt pour le petit modèle
                structure_result = self.cached_structure
                
                # Créer un prompt plus clair avec la structure exacte attendue
                regroup_prompt = f"""Tu es un assistant spécialisé dans la réorganisation de données de projets. Tu vas regrouper les informations de projets similaires ou liés.

TÂCHE:
Analyser la structure de données ci-dessous et retourner une version réorganisée sous EXACTEMENT le même format.

RÈGLES:
1. Regrouper les projets similaires ou liés en un seul projet
2. Ne pas modifier les informations des projets, seulement les regrouper
3. Maintenir la hiérarchie de projets, sous-projets, etc.
4. Conserver tous les événements à venir par service
5. Éviter les duplications d'informations/alertes/événements
6. IMPORTANT: Ton output DOIT contenir exactement la même structure que l'input, avec les mêmes clés au premier niveau ("projects", "upcoming_events", "metadata", "source_files")

STRUCTURE ATTENDUE:
{{
  "projects": {{
    "Nom_du_projet": {{
      "information": "texte",
      "critical": ["alerte 1", "alerte 2"],
      "small": ["alerte mineure 1", "alerte mineure 2"],
      "advancements": ["avancement 1", "avancement 2"]
    }},
    // ... autres projets
  }},
  "upcoming_events": {{
    "Service1": ["événement 1", "événement 2"],
    // ... autres services
  }},
  "metadata": {{ ... }},
  "source_files": [ ... ]
}}

STRUCTURE À ANALYSER:
{json.dumps(structure_result, indent=2, ensure_ascii=False)}

RENVOIE UNIQUEMENT LE JSON RÉORGANISÉ AVEC LA STRUCTURE COMPLÈTE."""
                
                try:
                    # Utiliser le petit modèle avec un contexte réduit
                    regrouped_structure = remove_tags_no_keep(self.small_model.invoke(regroup_prompt), '<think>', '</think>')
                    
                    # Log the raw response for debugging
                    log.info(f"Raw model response: {regrouped_structure}")
                    
                    # Use extract_json to handle any formatting issues
                    try:
                        regrouped_data = extract_json(regrouped_structure)
                        if not regrouped_data:
                            raise ValueError("No valid JSON could be extracted from the model's response")
                        
                        # Vérifier et corriger la structure si nécessaire
                        log.info(f"Extracted JSON structure keys: {regrouped_data.keys() if isinstance(regrouped_data, dict) else 'Not a dict'}")
                        
                        # Si le modèle n'a retourné que la partie "projects", recréer la structure complète
                        if isinstance(regrouped_data, dict) and "projects" not in regrouped_data:
                            # Vérifier si le modèle a directement retourné le contenu de "projects"
                            # (il y aurait des sous-dictionnaires avec information, critical, small, advancements)
                            has_project_structure = False
                            for key, value in regrouped_data.items():
                                if isinstance(value, dict) and any(k in value for k in ["information", "critical", "small", "advancements"]):
                                    has_project_structure = True
                                    break
                            
                            if has_project_structure:
                                log.info("Modèle a retourné directement le contenu de 'projects', reconstruction de la structure complète")
                                original_data = regrouped_data
                                regrouped_data = {
                                    "projects": original_data,
                                    "upcoming_events": structure_result.get("upcoming_events", {}),
                                    "metadata": structure_result.get("metadata", {}),
                                    "source_files": structure_result.get("source_files", [])
                                }
                            else:
                                log.error(f"Structure JSON invalide: {regrouped_data}")
                                raise ValueError("Invalid JSON structure: missing required fields and not a project structure")
                        
                        # Vérification finale de la structure
                        if not isinstance(regrouped_data, dict) or "projects" not in regrouped_data:
                            log.error(f"Structure JSON invalide après correction: {regrouped_data}")
                            raise ValueError("Invalid JSON structure: missing required fields after correction")
                        
                    except Exception as e:
                        log.error(f"JSON extraction/validation error: {str(e)}")
                        log.error(f"Invalid content: {regrouped_structure}")
                        response = f"Erreur lors de la réorganisation des données: impossible d'extraire un JSON valide. Détails: {str(e)}"
                        yield f"data: {json.dumps({'choices': [{'message': {'content': response}}]})}\n\n"
                        yield f"data: {json.dumps({'choices': [{'finish_reason': 'stop'}]})}\n\n"
                        self.last_response = response
                        return
                    
                    # Créer directement la présentation en utilisant update_table_with_project_data
                    from src.services.update_pttx_service import update_table_with_project_data
                    from pptx import Presentation
                    from pptx.util import Pt
                    
                    # Créer les dossiers nécessaires
                    output_regroup = os.path.join("OUTPUT", self.chat_id, "regrouped")
                    os.makedirs(output_regroup, exist_ok=True)
                    
                    # Nom du fichier de sortie avec timestamp
                    import datetime
                    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
                    output_file = os.path.join(output_regroup, f"regrouped_{timestamp}.pptx")
                    
                    # Créer une présentation vide ou utiliser un template
                    template_path = os.path.join("templates", "CRA_TEMPLATE_IA.pptx")
                    if os.path.exists(template_path):
                        prs = Presentation(template_path)
                    else:
                        prs = Presentation()
                        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
                        table_shape = slide.shapes.add_table(rows=10, cols=3, left=Pt(30), top=Pt(30), width=Pt(600), height=Pt(400))
                    
                    # Sauvegarder temporairement
                    temp_path = os.path.join(output_regroup, "temp.pptx")
                    prs.save(temp_path)
                    
                    # Mettre à jour la présentation avec les données regroupées
                    try:
                        # Utiliser le premier slide (index 0) et la première table (index 0)
                        updated_path = update_table_with_project_data(
                            temp_path, 
                            0, 
                            0, 
                            regrouped_data["projects"],
                            output_file,
                            regrouped_data.get("upcoming_events", {})
                        )
                        
                        # Supprimer le fichier temporaire
                        if os.path.exists(temp_path):
                            os.remove(temp_path)
                        
                        # Upload to OpenWebUI and get download URL
                        upload_result = self.download_file_openwebui(updated_path)
                        if "error" in upload_result:
                            response = f"Les informations des projets ont été regroupées avec succès, mais une erreur s'est produite lors de la génération du lien de téléchargement: {upload_result['error']}"
                        else:
                            response = f"Les informations des projets ont été regroupées avec succès.\n\n### URL de téléchargement:\n{upload_result.get('download_url', 'Non disponible')}"
                            # Save mappings after uploading the new file
                            self.save_file_mappings()
                    except Exception as e:
                        log.error(f"Error during PowerPoint generation: {str(e)}")
                        log.exception("PowerPoint generation error details:")
                        response = f"Erreur lors de la génération de la présentation PowerPoint: {str(e)}"
                    
                except Exception as e:
                    log.error(f"Error during regrouping: {str(e)}")
                    log.exception("Full error details:")
                    response = f"Erreur lors de la réorganisation des données: {str(e)}"
            
            yield f"data: {json.dumps({'choices': [{'message': {'content': response}}]})}\n\n"
            yield f"data: {json.dumps({'choices': [{'finish_reason': 'stop'}]})}\n\n"
            self.last_response = response
            return
        # Ajouter la dernière réponse au contexte si elle existe
        if user_message:
            user_message += f"\n\n *Last response generated :* {self.last_response}"
        else:
            # Afficher les commandes disponibles si aucune réponse précédente
            commands = """Les commandes sont les suivantes : 

/summarize [instructions] --> Affiche les résumés existants et demande confirmation avant d'en générer un nouveau. Vous pouvez ajouter des instructions spécifiques après la commande pour guider le résumé.
/structure --> Renvoie la structure des fichiers 
/clear [IDs] --> Nettoie tous les dossiers orphelins et supprime les fichiers associés dans OpenWebUI (préserve la conversation actuelle et éventuellement d'autres IDs spécifiés)
/generate --> Génère tout le pptx en fonction du texte ( /generate [Avancements de la semaine])
/merge --> Fusionne tous les fichiers pptx envoyés
/regroup --> Regroupe les informations des projets similaires ou liés

NOTE: Le système a été amélioré pour éviter la duplication d'informations. Les alertes (en couleur) apparaissent maintenant uniquement dans leur section correspondante et ne sont plus répétées dans les informations générales.
            """
            self.last_response = commands
            yield f"data: {json.dumps({'choices': [{'message': {'content': commands}}]})}\n\n"
            yield f"data: {json.dumps({'choices': [{'finish_reason': 'stop'}]})}\n\n"
            return
        
        # Initialiser le contenu cumulatif
        cumulative_content = ""
        user_message = self.system_prompt + "\n\n" + user_message
        
        try:
            # Format standard OpenAI-like qui est attendu par OpenWebUI
            # Premier message pour initialiser le stream
            yield f"data: {json.dumps({'choices': [{'delta': {'role': 'assistant'}}]})}\n\n"
            
            # Streamer la réponse depuis le modèle
            for chunk in self.streaming_model.stream(user_message):
                if isinstance(chunk, str):
                    content_delta = chunk
                else:
                    content_delta = chunk.content if hasattr(chunk, 'content') else str(chunk)
                
                # Nettoyer le contenu pour éviter les problèmes de formatage
                content_delta = content_delta.replace('\r', '')
                
                # Ajouter au contenu cumulatif
                cumulative_content += content_delta
                
                # Envoi de l'événement au client si un émetteur est disponible
                if __event_emitter__:
                    __event_emitter__({"type": "content_delta", "delta": content_delta})
                
                # Format compatible avec le standard OpenAI utilisé par OpenWebUI
                delta_response = {
                    "choices": [
                        {
                            "delta": {"content": content_delta}
                        }
                    ]
                }
                
                # Yield en format SSE (Server-Sent Events)
                yield f"data: {json.dumps(delta_response)}\n\n"
                
            # Message de fin spécifique
            yield f"data: {json.dumps({'choices': [{'delta': {}, 'finish_reason': 'stop'}]})}\n\n"
            yield f"data: [DONE]\n\n"  # Signal de fin standard OpenAI
            
        except Exception as e:
            error_message = f"Erreur lors du streaming de la réponse: {str(e)}"
            if __event_emitter__:
                __event_emitter__({"type": "error", "error": error_message})
            yield f"data: {json.dumps({'error': error_message})}\n\n"
            yield f"data: [DONE]\n\n"  # Même en cas d'erreur, on ferme proprement
            return
        
        self.last_response = cumulative_content

    def delete_file(self, file_path, update_mapping=True):
        """
        Supprime un fichier spécifique et met à jour le mapping si nécessaire.
        
        Args:
            file_path (str): Chemin du fichier à supprimer
            update_mapping (bool): Si True, met à jour le mapping des fichiers
            
        Returns:
            dict: Résultat de l'opération
        """
        try:
            abs_path = os.path.abspath(file_path)
            log.info(f"Deleting file: {abs_path}")
            
            if not os.path.exists(abs_path):
                log.warning(f"File not found: {abs_path}")
                return {"error": "Fichier non trouvé"}
            
            # Delete the file
            os.remove(abs_path)
            log.info(f"File deleted successfully: {abs_path}")
            
            # Update mapping if requested
            if update_mapping and abs_path in self.file_id_mapping:
                file_id = self.file_id_mapping[abs_path]
                del self.file_id_mapping[abs_path]
                log.info(f"Removed mapping for file {abs_path} with ID {file_id}")
                self.save_file_mappings()
            
            return {"message": f"Fichier {os.path.basename(file_path)} supprimé avec succès"}
        except Exception as e:
            log.error(f"Error deleting file {file_path}: {str(e)}")
            return {"error": f"Erreur lors de la suppression du fichier: {str(e)}"}

    def cleanup_orphaned_mappings(self):
        """
        Nettoie les mappings qui pointent vers des fichiers qui n'existent plus.
        
        Returns:
            int: Nombre de mappings supprimés
        """
        orphaned_mappings = []
        
        # Find all mappings pointing to files that no longer exist
        for file_path, file_id in list(self.file_id_mapping.items()):
            if not os.path.exists(file_path):
                orphaned_mappings.append((file_path, file_id))
        
        # Remove orphaned mappings
        removed_count = 0
        for file_path, file_id in orphaned_mappings:
            log.info(f"Removing orphaned mapping: {file_path} -> {file_id}")
            del self.file_id_mapping[file_path]
            removed_count += 1
        
        # Save updated mappings if any were removed
        if removed_count > 0:
            log.info(f"Removed {removed_count} orphaned mappings")
            self.save_file_mappings()
        else:
            log.info("No orphaned mappings found")
        
        return removed_count

    def force_cleanup_old_folders(self, exclude_chat_ids=None):
        """
        Méthode simplifiée pour nettoyer les dossiers qui ne correspondent pas aux IDs de chat exclus.
        Cette méthode est utilisée comme dernier recours pour nettoyer les anciens dossiers.
        
        Args:
            exclude_chat_ids (list): Liste des IDs de chat à conserver
            
        Returns:
            dict: Résultats du nettoyage
        """
        if exclude_chat_ids is None:
            exclude_chat_ids = []
            
        # Toujours préserver le chat_id actuel
        if self.chat_id and self.chat_id not in exclude_chat_ids:
            exclude_chat_ids.append(self.chat_id)
            
        log.info(f"Force cleaning folders, excluding: {exclude_chat_ids}")
        
        # Obtenir également la liste des chats actifs pour comparer et ajouter des mesures de sécurité
        active_chats = self.get_active_conversation_ids()
        log.info(f"Active chats from database: {active_chats}")
        
        # Ajouter à la liste d'exclusion les chats actifs qui n'y sont pas déjà
        for chat_id in active_chats:
            if chat_id not in exclude_chat_ids:
                log.info(f"Adding active chat to exclusion list: {chat_id}")
                exclude_chat_ids.append(chat_id)
                
        log.info(f"Final exclusion list: {exclude_chat_ids}")
        
        deleted_folders = 0
        deleted_files = 0
        deleted_webui_files = 0
        
        # Lister tous les dossiers à nettoyer
        folders_to_clean = []
        
        # Trouver les dossiers dans UPLOAD_FOLDER
        if os.path.exists(UPLOAD_FOLDER):
            for folder_name in os.listdir(UPLOAD_FOLDER):
                folder_path = os.path.join(UPLOAD_FOLDER, folder_name)
                if os.path.isdir(folder_path) and folder_name not in exclude_chat_ids:
                    folders_to_clean.append(folder_name)
        
        # Trouver les dossiers dans OUTPUT_FOLDER qui ne sont pas déjà dans la liste
        if os.path.exists(OUTPUT_FOLDER):
            for folder_name in os.listdir(OUTPUT_FOLDER):
                folder_path = os.path.join(OUTPUT_FOLDER, folder_name)
                if os.path.isdir(folder_path) and folder_name not in exclude_chat_ids and folder_name not in folders_to_clean:
                    folders_to_clean.append(folder_name)
        
        log.info(f"Dossiers à nettoyer: {folders_to_clean}")
        
        # Mesure de sécurité - Si tous les dossiers sont exclus, ne pas continuer
        if not folders_to_clean:
            log.info("Aucun dossier à nettoyer après exclusion des chats actifs.")
            return {
                "status": "success",
                "message": "Aucun dossier à nettoyer",
                "deleted_folders": 0,
                "deleted_files": 0,
                "deleted_webui_files": 0,
                "preserved_chats": exclude_chat_ids,
                "cleaned_chats": []
            }
        
        # Pour chaque dossier à nettoyer
        for chat_id in folders_to_clean:
            # Mesure de sécurité - Vérifier si le dossier contient des fichiers récents (moins de 1 jour)
            # Si c'est le cas, le considérer comme actif et l'ignorer
            upload_folder_path = os.path.join(UPLOAD_FOLDER, chat_id)
            output_folder_path = os.path.join(OUTPUT_FOLDER, chat_id)
            
            is_recent = False
            try:
                # Vérifier si des fichiers ont été modifiés récemment (moins de 24h)
                current_time = time.time()
                one_day_ago = current_time - (24 * 60 * 60)  # 24 heures en secondes
                
                # Vérifier dans le dossier UPLOAD
                if os.path.exists(upload_folder_path):
                    for filename in os.listdir(upload_folder_path):
                        file_path = os.path.join(upload_folder_path, filename)
                        if os.path.isfile(file_path):
                            mod_time = os.path.getmtime(file_path)
                            if mod_time > one_day_ago:
                                log.warning(f"Récent fichier trouvé dans {upload_folder_path}: {filename}. Préservation du dossier.")
                                is_recent = True
                                break
                
                # Vérifier dans le dossier OUTPUT
                if not is_recent and os.path.exists(output_folder_path):
                    for filename in os.listdir(output_folder_path):
                        file_path = os.path.join(output_folder_path, filename)
                        if os.path.isfile(file_path):
                            mod_time = os.path.getmtime(file_path)
                            if mod_time > one_day_ago:
                                log.warning(f"Récent fichier trouvé dans {output_folder_path}: {filename}. Préservation du dossier.")
                                is_recent = True
                                break
            except Exception as e:
                log.error(f"Erreur lors de la vérification des fichiers récents: {str(e)}")
                # En cas d'erreur, considérer le dossier comme récent par sécurité
                is_recent = True
            
            if is_recent:
                log.info(f"Préservation du chat {chat_id} car des fichiers récents y ont été trouvés")
                continue
            
            # 1. D'abord, supprimer les fichiers dans OpenWebUI
            log.info(f"Suppression des fichiers OpenWebUI pour le chat: {chat_id}")
            webui_result = self.delete_openwebui_files_for_chat(chat_id)
            deleted_webui_files += webui_result.get("deleted_count", 0)
            log.info(f"Résultat de suppression OpenWebUI: {webui_result}")
            
            # 2. Ensuite, supprimer le dossier dans UPLOAD_FOLDER
            if os.path.exists(upload_folder_path):
                log.info(f"Suppression du dossier: {upload_folder_path}")
                try:
                    file_count = len([f for f in os.listdir(upload_folder_path) 
                                    if os.path.isfile(os.path.join(upload_folder_path, f))])
                    deleted_files += file_count
                    shutil.rmtree(upload_folder_path, ignore_errors=True)
                    deleted_folders += 1
                    log.info(f"Dossier supprimé: {upload_folder_path}")
                except Exception as e:
                    log.error(f"Erreur lors de la suppression du dossier {upload_folder_path}: {str(e)}")
            
            # 3. Supprimer le dossier dans OUTPUT_FOLDER
            if os.path.exists(output_folder_path):
                log.info(f"Suppression du dossier: {output_folder_path}")
                try:
                    file_count = len([f for f in os.listdir(output_folder_path) 
                                    if os.path.isfile(os.path.join(output_folder_path, f))])
                    deleted_files += file_count
                    shutil.rmtree(output_folder_path, ignore_errors=True)
                    deleted_folders += 1
                    log.info(f"Dossier supprimé: {output_folder_path}")
                except Exception as e:
                    log.error(f"Erreur lors de la suppression du dossier {output_folder_path}: {str(e)}")
            
            # 4. Supprimer le fichier de mapping
            mapping_file = os.path.join(MAPPINGS_FOLDER, f"{chat_id}_file_mappings.json")
            if os.path.exists(mapping_file):
                log.info(f"Suppression du fichier de mapping: {mapping_file}")
                try:
                    os.remove(mapping_file)
                    log.info(f"Fichier de mapping supprimé: {mapping_file}")
                except Exception as e:
                    log.error(f"Erreur lors de la suppression du fichier de mapping {mapping_file}: {str(e)}")
            
        result = {
            "status": "success",
            "deleted_folders": deleted_folders,
            "deleted_files": deleted_files,
            "deleted_webui_files": deleted_webui_files,
            "preserved_chats": exclude_chat_ids,
            "cleaned_chats": folders_to_clean
        }
        
        log.info(f"Résultats du nettoyage forcé: {result}")
        return result

    def delete_file_from_openwebui(self, file_id, active_files_mapping=None):
        """
        Supprime un fichier dans OpenWebUI via l'API.
        
        Args:
            file_id (str): L'ID du fichier à supprimer
            active_files_mapping (dict): Dictionnaire contenant tous les mappings des conversations actives
            
        Returns:
            bool: True si la suppression a réussi, False sinon
        """
        try:
            if not file_id:
                log.warning("Tentative de suppression d'un fichier sans ID")
                return False
                
            # Vérifier si le fichier est utilisé par d'autres conversations actives
            if active_files_mapping and file_id in active_files_mapping:
                log.info(f"Le fichier {file_id} est toujours utilisé par d'autres conversations actives. Pas de suppression.")
                return False
                
            url = f"{self.openwebui_api}files/{file_id}"
            headers = {
                "accept": "application/json",
                "Authorization": f"Bearer {self.openwebui_api_key}"
            }
            
            log.info(f"Suppression du fichier OpenWebUI avec ID: {file_id}")
            response = requests.delete(url, headers=headers)
            
            if response.status_code in [200, 204]:
                log.info(f"Fichier {file_id} supprimé avec succès dans OpenWebUI")
                return True
            else:
                log.error(f"Échec de suppression du fichier {file_id} dans OpenWebUI: {response.status_code} - {response.text}")
                return False
        except Exception as e:
            log.error(f"Erreur lors de la suppression du fichier {file_id} dans OpenWebUI: {str(e)}")
            return False
            
    def get_all_active_files_mapping(self):
        """
        Récupère tous les fichiers utilisés par des conversations actives.
        
        Returns:
            dict: Dictionnaire des fichiers utilisés {file_id: [chat_ids]}
        """
        active_chats = self.get_active_conversation_ids()
        all_files = {}  # {file_id: [list of chat_ids]}
        
        # Parcourir tous les fichiers de mapping
        if os.path.exists(MAPPINGS_FOLDER):
            try:
                for filename in os.listdir(MAPPINGS_FOLDER):
                    if filename.endswith("_file_mappings.json"):
                        chat_id = filename.split("_file_mappings.json")[0]
                        
                        # Vérifier seulement les mappings des conversations actives
                        if chat_id in active_chats:
                            mapping_file = os.path.join(MAPPINGS_FOLDER, filename)
                            if os.path.exists(mapping_file):
                                try:
                                    with open(mapping_file, 'r') as f:
                                        try:
                                            chat_mappings = json.load(f)
                                            # Inverser le mapping pour avoir file_id: [chat_ids]
                                            for file_path, file_id in chat_mappings.items():
                                                if file_id not in all_files:
                                                    all_files[file_id] = []
                                                if chat_id not in all_files[file_id]:
                                                    all_files[file_id].append(chat_id)
                                        except json.JSONDecodeError:
                                            log.error(f"Erreur de décodage JSON pour {mapping_file}")
                                except Exception as e:
                                    log.error(f"Erreur de lecture du fichier {mapping_file}: {str(e)}")
            except Exception as e:
                log.error(f"Erreur lors du parcours du dossier des mappings: {str(e)}")
        
        log.info(f"Trouvé {len(all_files)} fichiers utilisés par des conversations actives")
        return all_files

    def delete_openwebui_files_for_chat(self, chat_id):
        """
        Supprime tous les fichiers OpenWebUI associés à un chat ID spécifique.
        Vérifie d'abord que les fichiers ne sont pas utilisés par d'autres conversations actives.
        
        Args:
            chat_id (str): L'ID du chat dont les fichiers doivent être supprimés
            
        Returns:
            dict: Résultats de l'opération
        """
        if not chat_id:
            log.warning("Tentative de suppression des fichiers sans ID de chat")
            return {"status": "error", "message": "ID de chat non spécifié", "deleted_count": 0}
            
        log.info(f"Suppression des fichiers OpenWebUI pour le chat: {chat_id}")
        
        # D'abord, obtenir tous les fichiers utilisés par des conversations actives
        active_files_mapping = self.get_all_active_files_mapping()
        
        # Charger le mapping pour ce chat
        mapping_file = os.path.join(MAPPINGS_FOLDER, f"{chat_id}_file_mappings.json")
        deleted_count = 0
        failed_count = 0
        skipped_count = 0
        file_ids = []
        
        try:
            # Vérifier si le fichier de mapping existe
            if os.path.exists(mapping_file):
                log.info(f"Fichier de mapping trouvé: {mapping_file}")
                try:
                    with open(mapping_file, 'r') as f:
                        mappings = json.load(f)
                    
                    # Extraire tous les IDs de fichiers du mapping
                    file_ids = list(set(mappings.values()))
                    log.info(f"Trouvé {len(file_ids)} fichiers potentiels à supprimer: {file_ids}")
                    
                    # Supprimer les fichiers qui ne sont pas utilisés par d'autres conversations
                    for file_id in file_ids:
                        # Si le fichier est utilisé par d'autres conversations que celle-ci
                        if file_id in active_files_mapping and len(active_files_mapping[file_id]) > 1:
                            other_chats = [c for c in active_files_mapping[file_id] if c != chat_id]
                            if other_chats:
                                log.info(f"Le fichier {file_id} est utilisé par d'autres conversations actives: {other_chats}. Ignoré.")
                                skipped_count += 1
                                continue
                        
                        # Le fichier n'est pas utilisé par d'autres conversations actives
                        if self.delete_file_from_openwebui(file_id):
                            deleted_count += 1
                        else:
                            failed_count += 1
                except json.JSONDecodeError as e:
                    log.error(f"Erreur de décodage du fichier mapping {mapping_file}: {str(e)}")
                    return {"status": "error", "message": f"Erreur de décodage JSON: {str(e)}", "deleted_count": 0}
            else:
                log.info(f"Aucun fichier de mapping trouvé pour le chat {chat_id}")
            
            # AJOUT: Suppression directe des fichiers dans le dossier uploads
            try:
                from services import delete_matching_files_in_openwebui
                log.info(f"Tentative de suppression directe des fichiers dans le dossier uploads pour {chat_id}")
                
                # Définir les variables d'environnement pour la fonction
                os.environ["OPENWEBUI_UPLOADS"] = "open-webui/uploads"
                
                # Essayer d'abord avec le chemin relatif
                if not os.path.exists("open-webui/uploads"):
                    # Essayer avec des chemins absolus
                    possible_paths = [
                        os.path.abspath("./open-webui/uploads"),
                        os.path.abspath("../open-webui/uploads"),
                        os.path.abspath("/open-webui/uploads"),
                        os.path.join(os.getcwd(), "open-webui/uploads")
                    ]
                    
                    for path in possible_paths:
                        log.info(f"Tentative avec le chemin: {path}")
                        if os.path.exists(path):
                            os.environ["OPENWEBUI_UPLOADS"] = path
                            log.info(f"Dossier uploads trouvé à: {path}")
                            break
                    else:
                        log.warning(f"Le dossier open-webui/uploads n'existe pas dans aucun des chemins essayés.")
                        log.info(f"Répertoire courant: {os.getcwd()}")
                        log.info(f"Contenu du répertoire courant: {os.listdir('.')}")
                
                # Appeler directement la fonction
                deleted_files = delete_matching_files_in_openwebui(chat_id)
                if deleted_files:
                    log.info(f"Suppression directe réussie: {len(deleted_files)} fichiers supprimés")
                    log.info(f"Fichiers supprimés: {deleted_files}")
                    deleted_count += len(deleted_files)
                else:
                    log.warning("Aucun fichier n'a pu être supprimé directement")
            except Exception as e:
                log.error(f"Erreur lors de la suppression directe des fichiers: {str(e)}")
        except Exception as e:
            log.error(f"Erreur lors de la suppression des fichiers OpenWebUI pour le chat {chat_id}: {str(e)}")
            return {
                "status": "error", 
                "message": f"Erreur: {str(e)}", 
                "deleted_count": deleted_count,
                "failed_count": failed_count,
                "skipped_count": skipped_count
            }
            
        return {
            "status": "success",
            "message": f"Supprimé {deleted_count} fichiers, {skipped_count} fichiers ignorés car encore utilisés, échec pour {failed_count} fichiers",
            "deleted_count": deleted_count,
            "failed_count": failed_count,
            "skipped_count": skipped_count,
            "file_ids": file_ids
        }

pipeline = Pipeline()

if __name__ == "__main__":
    summarize_ppt("1040706a-776f-4233-b823-b49658dc42dd")