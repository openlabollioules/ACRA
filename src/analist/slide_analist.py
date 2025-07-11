import os
from pptx import Presentation

# ---- Test for Color identification inside pptx ----

def get_run_color_tuple(run):
    """
    Retourne un tuple (R, G, B) pour la couleur du run s'il est accessible,
    sinon retourne None.
    """
    if run.font.color is None:
        return None
    try:
        rgb = run.font.color.rgb
        if rgb is None:
            return None
        # rgb est un objet de type RGBColor, qui se comporte comme une séquence
        return (rgb[0], rgb[1], rgb[2])
    except AttributeError:
        # Si la couleur est de type _SchemeColor ou inaccessible, on considère la couleur comme par défaut.
        return None

def is_default_color(color_tuple):
    """
    Considère qu'une couleur est par défaut si elle est None,
    ou si elle est noire (0,0,0) ou blanche (255,255,255).
    """
    if color_tuple is None:
        return True
    return color_tuple in [(0, 0, 0), (255, 255, 255)]

def process_text_frame(text_frame):
    """
    Concatène le texte de chaque paragraphe en insérant des balises de couleur
    pour les runs dont la couleur n'est pas noire ni blanche.
    
    Exemple de sortie :
    "le début de mon texte.. <rgb=255 0 0 >Ma partie en rouge...<rgb=255 0 0 > le reste de mon texte.."
    """
    result = ""
    for paragraph in text_frame.paragraphs:
        para_text = ""
        for run in paragraph.runs:
            text = run.text  # Conserve les espaces tels quels
            color_tuple = get_run_color_tuple(run)
            if not is_default_color(color_tuple):
                r, g, b = color_tuple
                # On insère une balise avant et après le run coloré
                para_text += f"<rgb={r} {g} {b} >{text}<rgb={r} {g} {b} >"
            else:
                para_text += text
        result += para_text + "\n"
    return result

# GPT AH CODE 
def analyze_presentation_with_colors(file_path="./pptx_folder/CRA_SERVICE_CYBER.pptx"):
    """
    Analyzes a PowerPoint file and returns a structured dictionary containing text with color tags
    and other elements like tables, images, and charts.
    """
    prs = Presentation(file_path)
    presentation_data = {
        "total_slides": len(prs.slides),
        "slides": []
    }

    for slide_index, slide in enumerate(prs.slides):
        slide_data = {
            "slide_number": slide_index + 1,
            "shapes": []
        }

        for shape_index, shape in enumerate(slide.shapes):
            shape_data = {
                "index": shape_index,
                "type": type(shape).__name__
            }

            # Handle text frames
            if shape.has_text_frame:
                shape_data["text"] = process_text_frame(shape.text_frame).strip()

            # Handle tables
            elif shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        if cell.text_frame:
                            cell_text = process_text_frame(cell.text_frame).strip()
                            row_data.append(cell_text)
                        else:
                            row_data.append("")
                    table_data.append(row_data)
                shape_data["table"] = table_data

            # Handle images
            elif shape.shape_type == 13:
                shape_data["is_image"] = True

            # Handle charts
            elif hasattr(shape, "has_chart") and shape.has_chart:
                chart_data = {
                    "type": str(shape.chart.chart_type),
                    "series": []
                }
                for series in shape.chart.plots[0].series:
                    chart_data["series"].append({
                        "name": series.name,
                        "values": [pt for pt in series.values]
                    })
                shape_data["chart"] = chart_data

            slide_data["shapes"].append(shape_data)

        presentation_data["slides"].append(slide_data)
    
    print(presentation_data)

    return presentation_data

if __name__ == "__main__":
    # os.chdir("")
    print(os.getcwd())
    analyze_presentation_with_colors()