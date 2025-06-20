from pptx import Presentation
import re
import json
from typing import Dict, List, Tuple, Optional

def is_underlined(run):
    """
    Check if a text run is underlined.
    """
    # Check if run has font attribute directly
    if hasattr(run, 'font') and hasattr(run.font, 'underline'):
        return run.font.underline
    
    # Check if run has parent with font attribute
    if hasattr(run, '_parent') and hasattr(run._parent, 'font') and hasattr(run._parent.font, 'underline'):
        return run._parent.font.underline
    
    return False

def is_bold(run):
    """
    Check if a text run is bold.
    """
    # Check if run has font attribute directly
    if hasattr(run, 'font'):
        # Check for bold in font name or bold attribute
        if (hasattr(run.font, 'name') and run.font.name and "bold" in run.font.name.lower()) or \
           (hasattr(run.font, 'bold') and run.font.bold):
            return True
    
    # Check if run has parent with font attribute
    if hasattr(run, '_parent') and hasattr(run._parent, 'font'):
        # Check for bold in parent font name or bold attribute
        if (hasattr(run._parent.font, 'name') and run._parent.font.name and "bold" in run._parent.font.name.lower()) or \
           (hasattr(run._parent.font, 'bold') and run._parent.font.bold):
            return True
    
    return False

def get_rgb_color(run):
    """
    Get the RGB color of a text run.
    Returns tuple (R, G, B) or None if color is not accessible.
    """
    # Try to get color from run's font
    if hasattr(run, 'font') and hasattr(run.font, 'color') and run.font.color is not None:
        if hasattr(run.font.color, 'rgb') and run.font.color.rgb is not None:
            return tuple(run.font.color.rgb)
    
    # Try to get color from run's parent font
    if hasattr(run, '_parent') and hasattr(run._parent, 'font') and \
       hasattr(run._parent.font, 'color') and run._parent.font.color is not None:
        if hasattr(run._parent.font.color, 'rgb') and run._parent.font.color.rgb is not None:
            return tuple(run._parent.font.color.rgb)
    
    return None

def identify_color_type(color_tuple: Tuple[int, int, int]) -> str:
    """
    Identify color type based on RGB values.
    - Green: big advancement
    - Orange: small alert
    - Red: critical alert
    """
    if color_tuple is None:
        return "normal"
    
    r, g, b = color_tuple
    
    # Simple heuristic for color identification
    if g > max(r, b) + 50:  # Green is dominant
        return "advancement"
    elif r > g + 50 and g > b + 50:  # Orange-ish
        return "small_alert"
    elif r > max(g, b) + 50:  # Red is dominant
        return "critical_alert"
    else:
        return "normal"

def extract_title_from_slide(slide) -> str:
    """
    Extract title from a text field in the slide.
    Returns the title text.
    """
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame:
            # Assuming the first text field with content is the title
            if shape.text_frame.text.strip():
                return shape.text_frame.text.strip()
    return "Untitled"

def extract_table_data_from_slide(slide) -> List[Dict]:
    """
    Extract table data from a slide, focusing on tables with 3 columns:
    1. Project name
    2. Project information
    3. Upcoming events
    Returns a list of rows with text and formatting information.
    """
    results = []
    table_count = 0
    
    for shape_idx, shape in enumerate(slide.shapes):
        print(f"Shape {shape_idx}: Type {type(shape).__name__}")
        
        if hasattr(shape, 'has_table'):
            print(f"Shape {shape_idx} has_table attribute: {shape.has_table}")
        
        if shape.has_table:
            table_count += 1
            table = shape.table
            print(f"Processing table {table_count} with {len(table.rows)} rows and {len(table.columns)} columns")
            
            # Verify that we have the expected table structure - at least 3 columns
            if len(table.columns) < 3:
                print(f"WARNING: Table does not have 3 columns (found {len(table.columns)}). Skipping.")
                continue
            
            # Process each row in the table
            row_processed = 0
            for row_idx, row in enumerate(table.rows):
                # Skip header row if it exists (optional)
                # You can uncomment this if your table has a header row to skip
                if row_idx == 0:
                    print(f"Skipping header row (row 0)")
                    continue
                
                row_data = []
                has_content = False
                
                # We only care about the 3 columns we expect - but avoid slicing
                for col_idx, cell in enumerate(row.cells):
                    # Only process the first 3 columns
                    if col_idx >= 3:
                        break
                        
                    cell_text = cell.text.strip() if hasattr(cell, 'text') else ""
                    print(f"Row {row_idx}, Column {col_idx}: Text length {len(cell_text)}")
                        
                    # Vérifier si la cellule a du contenu
                    if hasattr(cell, 'text_frame'):
                        cell_data = {
                            "text": "",
                            "paragraphs": [],
                            "column_index": col_idx  # Track which column this is
                        }
                        
                        # Traiter chaque paragraphe dans la cellule
                        for para_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                            para_text = paragraph.text.strip()
                            print(f"  Paragraph {para_idx}: Text length {len(para_text)}")
                            
                            para_data = {
                                "text": "",
                                "runs": []
                            }
                            
                            # Traiter chaque run dans le paragraphe
                            for run_idx, run in enumerate(paragraph.runs):
                                run_text = run.text
                                if run_text.strip():  # ignorer les runs vides
                                    color = get_rgb_color(run)
                                    color_type = identify_color_type(color)
                                    print(f"    Run {run_idx}: Text '{run_text[:20]}...' Color type: {color_type}")
                                    
                                    para_data["text"] += run_text
                                    para_data["runs"].append({
                                        "text": run_text,
                                        "color": color,
                                        "color_type": color_type
                                    })
                                    has_content = True
                            
                            # Ajouter le paragraphe seulement s'il contient du texte
                            if para_data["text"].strip():
                                cell_data["text"] += para_data["text"] + "\n"
                                cell_data["paragraphs"].append(para_data)
                            
                        cell_data["text"] = cell_data["text"].strip()
                        
                        # For column 2 (upcoming events), verify it's really an upcoming event
                        if col_idx == 2:
                            # Check if the text contains indicators of future events
                            text_lower = cell_data["text"].lower()
                            event_indicators = [
                                "événement", "evenement", "à venir", "a venir", "prochain", 
                                "semaine prochaine", "mois prochain", "futur", "prévu", "prevu",
                                "sera", "planning", "calendrier", "agenda", "rendez-vous", "rendez vous"
                            ]
                            
                            is_upcoming_event = any(indicator in text_lower for indicator in event_indicators)
                            
                            if not is_upcoming_event and cell_data["text"]:
                                print(f"Column 2 text doesn't appear to be an upcoming event: '{cell_data['text'][:30]}...'")
                                # Add a flag to indicate this might not be an upcoming event
                                cell_data["is_upcoming_event"] = False
                            else:
                                cell_data["is_upcoming_event"] = True
                        
                        row_data.append(cell_data)
                        if cell_data["text"]:
                            has_content = True
                    else:
                        # Ajouter une cellule vide avec l'index de colonne approprié
                        row_data.append({"text": "", "paragraphs": [], "column_index": col_idx})
                
                # Ajouter cette ligne aux résultats seulement si elle contient au moins des données dans la colonne 0 ou 1
                if any(cell.get("text", "").strip() for cell in row_data if cell.get("column_index") in [0, 1]):
                    results.append(row_data)
                    row_processed += 1
                    print(f"Row {row_idx} added to results")
                else:
                    print(f"Row {row_idx} skipped (no content in columns 0 or 1)")
            
            print(f"Processed {row_processed} rows from table {table_count}")
    
    print(f"Total tables found: {table_count}, Total rows extracted: {len(results)}")
    return results

def extract_projects_from_table_data(table_data: List[Dict], title: str) -> Dict[str, Dict]:
    """
    Extract project information from processed table data.
    Assuming the table has 3 columns:
    1. Project name (column_index=0)
    2. Project information with colored text for alerts (column_index=1)
    3. Upcoming events (column_index=2)
    Returns a dictionary with project names as keys and their information as values.
    
    Automatically detects project hierarchies based on name patterns:
    - Simple naming: "Project Subproject" -> {Project: {Subproject: {...}}}
    - Parenthesis naming: "Project (Subproject)" -> {Project: {Subproject: {...}}}
    - Multi-level naming: "Main Sub (Detail)" -> {Main: {Sub: {Detail: {...}}}}
    
    Note: upcoming_events from column 2 are collected separately and not stored at the project level.
    They will be used to populate the upcoming_events by service at a higher level.
    """
    # Initialize with a multi-level structure for projects hierarchy
    projects = {}
    collected_upcoming_events = []
    
    # Store raw project data first to analyze hierarchy later
    raw_projects = {}
    
    # Process each row in the table data to collect raw project data
    for row in table_data:
        # Get cells by column index
        project_name_cell = next((cell for cell in row if cell.get("column_index") == 0), {})
        project_info_cell = next((cell for cell in row if cell.get("column_index") == 1), {})
        events_cell = next((cell for cell in row if cell.get("column_index") == 2), {})
        
        # Extract project name from column 0
        full_project_name = project_name_cell.get("text", "").strip()
        if not full_project_name:
            continue
            
        # Clean and normalize the project name for comparison
        normalized_name = full_project_name.lower().strip()
        
        # Store raw data with the original name
        raw_projects[full_project_name] = {
            "normalized_name": normalized_name,
            "information": "",
            "critical": [],
            "small": [],
            "advancements": []
            # Note: upcoming_events n'est plus stocké au niveau du projet
        }
        
        # Process project information from column 1
        project_information = ""
        
        for paragraph in project_info_cell.get("paragraphs", []):
            # Track the original paragraph text
            paragraph_text = paragraph.get("text", "")
            
            # Keep the full paragraph text for information field
            # We'll also identify colored sections for alerts
            project_information += paragraph_text + "\n"
            
            # Process runs to extract colored alerts
            for run in paragraph.get("runs", []):
                run_text = run["text"]
                
                # Add text to appropriate category based on color
                if run["color_type"] == "advancement":
                    if run_text not in raw_projects[full_project_name]["advancements"]:
                        raw_projects[full_project_name]["advancements"].append(run_text)
                elif run["color_type"] == "small_alert":
                    if run_text not in raw_projects[full_project_name]["small"]:
                        raw_projects[full_project_name]["small"].append(run_text)
                elif run["color_type"] == "critical_alert":
                    if run_text not in raw_projects[full_project_name]["critical"]:
                        raw_projects[full_project_name]["critical"].append(run_text)
        
        # Set the information text
        raw_projects[full_project_name]["information"] = project_information.strip()
        
        # Process upcoming events from column 2 - collect them pour les remonter au niveau supérieur
        events_text = events_cell.get("text", "").strip()
        
        # Only add to upcoming events if it's verified as an actual upcoming event
        is_upcoming_event = events_cell.get("is_upcoming_event", True)  # Default to True for backward compatibility
        
        if events_text and is_upcoming_event and events_text not in collected_upcoming_events:
            print(f"Adding verified upcoming event: {events_text[:30]}...")
            collected_upcoming_events.append(events_text)
        elif events_text and not is_upcoming_event:
            print(f"Skipping text that's not an upcoming event: {events_text[:30]}...")
            # Instead of adding to upcoming events, we could add it to general information
            if raw_projects[full_project_name]["information"]:
                raw_projects[full_project_name]["information"] += "\n" + events_text
            else:
                raw_projects[full_project_name]["information"] = events_text
    
    # Function to extract hierarchy from project name
    def extract_hierarchy(name):
        # Try to match patterns like "Main Sub (Detail)" or "Main Sub Detail"
        
        # First check for parenthesis format: "Project (Subproject)"
        parenthesis_match = re.search(r'(.*?)\s*\((.*?)\)', name)
        if parenthesis_match:
            main_part = parenthesis_match.group(1).strip()
            sub_part = parenthesis_match.group(2).strip()
            
            # Check if main_part itself contains spaces indicating further hierarchy
            main_parts = main_part.split(' ')
            if len(main_parts) > 1:
                # Take first word as top-level project
                top_level = main_parts[0].strip()
                # Rest as mid-level
                mid_level = ' '.join(main_parts[1:]).strip()
                return [top_level, mid_level, sub_part]
            else:
                return [main_part, sub_part]
        
        # No parenthesis, check for space-separated parts
        parts = name.split(' ')
        if len(parts) >= 2:
            # First word as main project, rest as subproject
            return [parts[0], ' '.join(parts[1:])]
        
        # No clear hierarchy, treat as single project
        return [name]
    
    # Build the project hierarchy
    for original_name, data in raw_projects.items():
        # Extract hierarchy levels from the project name
        hierarchy = extract_hierarchy(original_name)
        
        # Convert hierarchy to lowercase for case-insensitive matching
        hierarchy_lower = [level.lower() for level in hierarchy]
        
        # Build nested structure
        current_level = projects
        for i, level in enumerate(hierarchy):
            level_lower = hierarchy_lower[i]
            
            # Find existing key with case-insensitive match
            existing_key = None
            for key in current_level.keys():
                if key.lower() == level_lower:
                    existing_key = key
                    break
            
            # Use the original case from the first occurrence we saw
            actual_key = existing_key if existing_key else level
            
            if i == len(hierarchy) - 1:  # Last level - add the data
                if actual_key not in current_level:
                    current_level[actual_key] = data.copy()
                    # Remove the normalized_name from the final data
                    del current_level[actual_key]["normalized_name"]
                else:
                    # Merge with existing data
                    current_level[actual_key]["information"] += "\n" + data["information"] if current_level[actual_key]["information"] else data["information"]
                    current_level[actual_key]["critical"].extend(data["critical"])
                    current_level[actual_key]["small"].extend(data["small"])
                    current_level[actual_key]["advancements"].extend(data["advancements"])
            else:
                # Create intermediate level if it doesn't exist
                if actual_key not in current_level:
                    current_level[actual_key] = {}
                
                # Move to next level
                current_level = current_level[actual_key]
    
    # Add metadata for reference
    metadata = {
        "title": title,
        "collected_upcoming_events": collected_upcoming_events  # Stocker les événements collectés dans les métadonnées
    }
    
    return {
        "projects": projects,
        "metadata": metadata
    }

def extract_projects_from_presentation(file_path: str) -> Dict[str, Dict]:
    """
    Extract project information from a PowerPoint presentation.
    Focuses on the first slide with a title and a 3-column table.
    """
    try:
        print(f"Attempting to process PowerPoint file: {file_path}")
        prs = Presentation(file_path)
        print(f"Successfully loaded presentation with {len(prs.slides)} slides")
        
        # Process only the first slide as specified
        if len(prs.slides) > 0:
            slide = prs.slides[0]
            title = extract_title_from_slide(slide)
            print(f"Extracted title: {title}")
            
            # Check if there are any tables in the slide
            has_tables = any(shape.has_table for shape in slide.shapes)
            print(f"Slide has tables: {has_tables}")
            
            # Count shapes in the slide
            shape_count = len(slide.shapes)
            print(f"Slide has {shape_count} shapes")
            
            table_data = extract_table_data_from_slide(slide)
            print(f"Extracted table data with {len(table_data)} rows")
            
            if not table_data:
                print("WARNING: No table data was extracted from the slide")
                # Return empty structure rather than failing
                return {
                    "projects": {},
                    "metadata": {
                        "title": title,
                        "collected_upcoming_events": [],
                        "error": "No table data extracted from slide"
                    }
                }
            
            projects = extract_projects_from_table_data(table_data, title)
            print(f"Extracted projects: {len(projects.get('projects', {}))} top-level projects")
            return projects
        else:
            print("WARNING: Presentation has no slides")
            return {
                "projects": {},
                "metadata": {
                    "title": "No slides",
                    "collected_upcoming_events": [],
                    "error": "Presentation has no slides"
                }
            }
    except Exception as e:
        error_message = f"Error processing presentation: {e}"
        print(error_message)
        # Return empty structure with error info instead of empty dict
        return {
            "projects": {},
            "metadata": {
                "title": "Error",
                "collected_upcoming_events": [],
                "error": error_message
            }
        }

def format_projects_as_json(projects: Dict[str, Dict], output_file: Optional[str] = None) -> str:
    """
    Format project information as JSON and optionally save to a file.
    """
    json_data = json.dumps(projects, indent=2, ensure_ascii=False)
    
    if output_file:
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(json_data)
    
    return json_data

def extract_and_format_projects(file_path: str, output_file: Optional[str] = None) -> Dict[str, Dict]:
    """
    Main function to extract project information from a PowerPoint presentation
    and format it as JSON.
    """
    projects = extract_projects_from_presentation(file_path)
    
    if output_file:
        format_projects_as_json(projects, output_file)
    
    return projects

if __name__ == "__main__":
    import sys
    
    if len(sys.argv) > 1:
        pptx_file = sys.argv[1]
        output_file = sys.argv[2] if len(sys.argv) > 2 else None
        projects = extract_and_format_projects(pptx_file, output_file)
        
        if not output_file:
            print(json.dumps(projects, indent=2, ensure_ascii=False))
    else:
        print("Usage: python project_extractor.py <pptx_file> [output_json_file]") 