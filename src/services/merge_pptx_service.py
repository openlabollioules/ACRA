import os
from pathlib import Path
from pptx import Presentation
import logging
from copy import deepcopy
import six
import copy
import tempfile
import shutil
from lxml import etree

def merge_pptx_files(input_folder: str, output_path: str) -> str:
    """
    Merge all PowerPoint (.pptx) files in a folder into a single presentation.
    
    Args:
        input_folder: Path to the folder containing .pptx files
        output_path: Path where the merged presentation will be saved
    
    Returns:
        Path to the merged presentation file
    
    Raises:
        FileNotFoundError: If input_folder doesn't exist
        ValueError: If no .pptx files are found in the folder
    """
    # Check if input folder exists
    input_folder_path = Path(input_folder)
    if not input_folder_path.exists() or not input_folder_path.is_dir():
        raise FileNotFoundError(f"Input folder not found: {input_folder}")
    
    # Find all .pptx files
    pptx_files = list(input_folder_path.glob("*.pptx"))
    
    if not pptx_files:
        raise ValueError(f"No PowerPoint files found in {input_folder}")
    
    # Sort files by name for consistent order
    pptx_files = sorted(pptx_files)
    
    # Create output directory if it doesn't exist
    output_file_path = Path(output_path)
    os.makedirs(output_file_path.parent, exist_ok=True)
    
    # Create a base empty presentation with the target file
    base_prs = Presentation()
    base_prs.save(output_path)
    
    # Process each presentation
    for pptx_file in pptx_files:
        try:
            # For each presentation, copy slides one by one to the target
            copy_slides_between_presentations(str(pptx_file), output_path)
            logging.info(f"Added slides from {pptx_file}")
        except Exception as e:
            logging.error(f"Error processing {pptx_file}: {str(e)}")
    
    return output_path

def copy_slides_between_presentations(source_path, target_path):
    """
    Copy all slides from source presentation to target presentation.
    
    Args:
        source_path: Path to source presentation
        target_path: Path to target presentation
    """
    # Open source presentation
    source_prs = Presentation(source_path)
    
    # Process each slide in the source presentation
    for slide_index in range(len(source_prs.slides)):
        # Create a temporary presentation with just one slide
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_slide_path = os.path.join(temp_dir, "temp_slide.pptx")
            
            # Create a new presentation and copy one slide from source
            temp_prs = Presentation()
            
            # Get the source slide
            source_slide = source_prs.slides[slide_index]
            
            # Find the right layout in the temp presentation
            try:
                # Try to get a blank layout for maximum compatibility
                blank_layout = temp_prs.slide_layouts[6]  # Usually blank layout
            except:
                # Fallback to the last layout
                blank_layout = temp_prs.slide_layouts[-1]
            
            # Add a slide to the temp presentation
            temp_slide = temp_prs.slides.add_slide(blank_layout)
            
            # Copy all shapes from source to temp slide (excluding placeholders)
            for shape in source_slide.shapes:
                el = shape.element
                # Skip placeholder elements (template elements)
                if el.tag.endswith('ph'):
                    continue
                
                # Create a deep copy of the shape element
                try:
                    newel = copy.deepcopy(el)
                    # Insert the copied element into the temp slide
                    temp_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
                except Exception as e:
                    logging.warning(f"Failed to copy shape: {str(e)}")
            
            # Save the temp presentation with one slide
            temp_prs.save(temp_slide_path)
            
            # Now append this slide to the target presentation
            append_slide_to_presentation(temp_slide_path, target_path)

def append_slide_to_presentation(source_path, target_path):
    """
    Append the first slide from source presentation to the target presentation.
    
    Args:
        source_path: Path to source presentation with one slide
        target_path: Path to target presentation
    """
    # Open the target presentation
    target_prs = Presentation(target_path)
    
    # Open the source presentation
    source_prs = Presentation(source_path)
    
    # Get the first slide from source
    source_slide = source_prs.slides[0]
    
    # Find the right layout in the target presentation
    try:
        blank_layout = target_prs.slide_layouts[6]  # Usually blank layout
    except:
        blank_layout = target_prs.slide_layouts[-1]
    
    # Add a slide to the target presentation
    target_slide = target_prs.slides.add_slide(blank_layout)
    
    # Copy all shapes from source to target slide (excluding placeholders)
    for shape in source_slide.shapes:
        el = shape.element
        # Skip placeholder elements (template elements)
        if el.tag.endswith('ph'):
            continue
        
        # Create a deep copy of the shape element
        try:
            newel = copy.deepcopy(el)
            # Insert the copied element into the target slide
            target_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        except Exception as e:
            logging.warning(f"Failed to copy shape: {str(e)}")
    
    # Save the updated target presentation
    target_prs.save(target_path)

def copy_slide(pres_path, target_pres_path):
    """
    Copy a slide from one presentation file to another.
    
    Args:
        pres_path: Path to source presentation
        target_pres_path: Path to target presentation
        
    Returns:
        The new slide created in the target presentation
    """
    index = 0
    target_pres = Presentation()
    pres = Presentation(pres_path)
    copy_template = pres.slides[index]

    try:
        blank_slide_layout = pres.slide_layouts[12]
    except:
        blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts) - 1]
    
    copied_slide = target_pres.slides.add_slide(blank_slide_layout)

    # Filter out hidden template elements
    for shp in copy_template.shapes:
        el = shp.element
        # Skip placeholder elements (usually template elements)
        if el.tag.endswith('ph'):
            continue
        
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    target_pres.save(target_pres_path)

    return copied_slide

# if __name__ == "__main__":
    # Example usage
    # merge_pptx_files("./pptx_folder/1", "./OUTPUT/1/merged_presentation.pptx")