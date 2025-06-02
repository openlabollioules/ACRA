"""
Command Handler Service for ACRA
Centralized command processing for the pipeline
"""
import os
import json
import datetime
from typing import Dict, Any, Generator, List, Tuple, Optional
from OLLibrary.utils.log_service import get_logger
from OLLibrary.utils.json_service import extract_json
from config_pipeline import acra_config
from .file_manager import FileManager
from .model_manager import model_manager
from .cleanup_service import cleanup_orphaned_folders

# Imports for PowerPoint generation (potentially move to a dedicated service later)
from pptx import Presentation
from pptx.util import Pt, Inches
from src.services.update_pttx_service import update_table_with_project_data

log = get_logger(__name__)

class CommandHandler:
    """
    Centralized command handler for ACRA pipeline commands.
    Processes commands like /summarize, /structure, /clear, etc.
    """
    
    def __init__(self, file_manager: FileManager):
        self.file_manager = file_manager
        self.cached_structure: Optional[Dict[str, Any]] = None
        self.waiting_for_confirmation = False
        self.confirmation_command = ""
        self.confirmation_additional_info: Optional[str] = None
        self.last_response: Optional[str] = None
        self.system_prompt: str = ""
    
    def reset_state(self):
        """Reset command handler state"""
        self.cached_structure = None
        self.waiting_for_confirmation = False
        self.confirmation_command = ""
        self.confirmation_additional_info = None
        self.last_response = None
    
    def get_available_commands(self) -> str:
        """
        Get a list of available commands with their descriptions.
        
        Returns:
            str: Formatted string with available commands and descriptions
        """
        return """Available commands: 

/summarize [instructions] --> Displays existing summaries and asks for confirmation before generating a new one. You can add specific instructions after the command to guide the summary.
/structure --> Returns the structure of the files 
/clear [IDs] --> Cleans all orphaned folders and deletes associated files in OpenWebUI (preserves the current conversation and optionally specified IDs)
/generate --> Generates a complete PowerPoint based on text (/generate [Weekly progress])
/merge --> Merges all uploaded PowerPoint files
/regroup --> Regroups information from similar or related projects"""
    
    def handle_confirmation(self, message: str) -> Tuple[bool, str]:
        """
        Handle confirmation responses (yes/no) from the user.
        
        This method is called when the system is waiting for a confirmation
        response from the user, such as after asking if they want to generate
        a new summary when summaries already exist.
        
        Args:
            message (str): The user's response message
            
        Returns:
            Tuple[bool, str]: (handled, response_message)
                - handled: True if the confirmation was handled, False otherwise
                - response_message: The response message if handled, empty string otherwise
        """
        if not self.waiting_for_confirmation:
            return False, ""
        
        message_lower = message.lower()
        
        if message_lower in ["yes", "y", "oui", "o"]:
            self.waiting_for_confirmation = False
            
            if self.confirmation_command == "summarize":
                return True, self._execute_summarize(self.confirmation_additional_info)
            
        elif message_lower in ["no", "n", "non"]:
            self.waiting_for_confirmation = False
            return True, "Summary generation canceled."
        
        # Reset if we get any other input
        self.waiting_for_confirmation = False
        return False, ""
    
    def handle_summarize_command(self, message: str) -> str:
        """
        Handle the /summarize command, which generates a summary PowerPoint from uploaded files.
        
        This method:
        1. Extracts any additional information provided after the command
        2. Checks for existing summaries
        3. If summaries exist, asks for confirmation before generating a new one
        4. If no summaries exist, proceeds directly to generation
        
        Args:
            message (str): The user's message containing the /summarize command and optional additional info
            
        Returns:
            str: Response message with either existing summaries, confirmation request, or generation result
        """
        # Extract additional information after the command
        additional_info = None
        if " " in message:
            command_parts = message.split(" ", 1)
            if len(command_parts) > 1 and command_parts[1].strip():
                additional_info = command_parts[1].strip()
        
        # Get existing summaries
        existing_summaries = self.file_manager.get_existing_summaries()
        
        if existing_summaries:
            # If there are existing summaries, show them and ask for confirmation before generating a new one
            response = "Voici les résumés existants pour cette conversation:\n\n"

            for filename, url in existing_summaries:
                response += f"- {filename}: {url}\n"
            
            response += "\nVoulez-vous générer un nouveau résumé? (Oui/Non)"
            
            # Set state to wait for confirmation
            self.waiting_for_confirmation = True
            self.confirmation_command = "summarize"
            self.confirmation_additional_info = additional_info
            
            return response
        else:
            # No existing summaries, generate one directly
            return self._execute_summarize(additional_info)
    
    def _generate_summary_powerpoint(self, summarized_structure: Dict[str, Any], timestamp: str) -> str:
        """
        Generates a PowerPoint presentation from a summarized JSON structure.

        Args:
            summarized_structure (Dict[str, Any]): The LLM-summarized project data.
            timestamp (str): Timestamp for generating a unique filename.

        Returns:
            str: Absolute path to the generated PowerPoint file, or an error string.
        """
        try:
            chat_id = self.file_manager.chat_id
            if not chat_id:
                log.error("Cannot generate summary PowerPoint without chat_id.")
                return "error: Chat ID not set for summary generation."

            # Define output directory for summaries within the chat_id's output folder
            summary_output_dir = os.path.join(acra_config.get_conversation_output_folder(chat_id), "summaries")
            os.makedirs(summary_output_dir, exist_ok=True)

            output_filename = f"summary_{timestamp}.pptx"
            output_filepath = os.path.join(summary_output_dir, output_filename)
            temp_filepath = os.path.join(summary_output_dir, f"temp_summary_{timestamp}.pptx")

            log.info(f"Generating summary PowerPoint at: {output_filepath}")

            # Create presentation: Use template if available, otherwise a blank one
            if acra_config.template_path and os.path.exists(acra_config.template_path):
                log.info(f"Using template: {acra_config.template_path}")
                prs = Presentation(acra_config.template_path)
                # Ensure the template has at least one slide and a table placeholder, or adapt as needed.
                # This example assumes the first slide and first shape (if a table) is the target.
                # More robust template handling might be needed (e.g., named placeholders).
                if not prs.slides:
                    log.warning("Template has no slides. Adding a blank slide.")
                    prs.slides.add_slide(prs.slide_layouts[5]) # Fallback to a blank slide layout
            else:
                log.info("No valid template found or specified. Creating a blank presentation.")
                prs = Presentation()
                # Add a blank slide (layout 5 is typically blank)
                slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(slide_layout)
                # Add a table placeholder - dimensions might need adjustment
                left = top = Inches(1.0)
                width = Inches(8.0)
                height = Inches(5.5)
                # Add a table with a default size, update_table_with_project_data should handle actual content
                # The row/col count here is a placeholder; update_table_with_project_data will manage it.
                slide.shapes.add_table(2, 2, left, top, width, height) 

            prs.save(temp_filepath) # Save initial state (template or blank with table)

            # Extract projects and upcoming_events for update_table_with_project_data
            projects_data = summarized_structure.get("projects", {})
            upcoming_events_data = summarized_structure.get("upcoming_events", {})

            # Call update_table_with_project_data to populate the presentation
            # Assuming slide_index=0 and table_index=0 for simplicity.
            # This might need to be more dynamic if templates have specific structures.
            final_pptx_path = update_table_with_project_data(
                pptx_path=temp_filepath, 
                slide_index=0, 
                table_shape_index=0, 
                project_data=projects_data,
                output_path=output_filepath, # Final desired output path
                upcoming_events=upcoming_events_data
            )

            if os.path.exists(temp_filepath):
                os.remove(temp_filepath)

            if "error" in final_pptx_path.lower(): # Check if update_table_with_project_data returned an error string
                 log.error(f"Error from update_table_with_project_data: {final_pptx_path}")
                 return f"error: Failed to update PowerPoint table - {final_pptx_path}"
            
            log.info(f"Summary PowerPoint generated successfully: {final_pptx_path}")
            return final_pptx_path # Should be the same as output_filepath if successful

        except Exception as e:
            log.error(f"Error in _generate_summary_powerpoint: {str(e)}", exc_info=True)
            if os.path.exists(temp_filepath):
                try: os.remove(temp_filepath) # Clean up temp file on error
                except: pass
            return f"error: Exception generating summary PowerPoint - {str(e)}"

    def _execute_summarize(self, additional_info: Optional[str] = None) -> str:
        """
        Execute the summarization operation, generating a new summary PowerPoint.
        
        This method orchestrates the entire summarization pipeline, with two possible paths:
        1. API path: Uses a remote API endpoint to generate the summary structure, then creates a PowerPoint locally
        2. Non-API path: Uses core.summarize_ppt which handles both summarization and PowerPoint generation
        
        Args:
            additional_info (str, optional): Additional context information to include in the summary
            
        Returns:
            str: Response message with download URL or error information
        """
        try:
            from core import summarize_ppt
            
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            chat_id = self.file_manager.chat_id

            if not chat_id:
                log.error("Chat ID not set. Cannot execute summarize.")
                return "Erreur: Chat ID non défini."

            generated_pptx_path: Optional[str] = None # To store the path of the generated PPTX

            if acra_config.get("USE_API"):
                # === API-BASED WORKFLOW ===
                # In this path, we use an API to generate the summary structure, then create the PowerPoint locally
                log.info(f"Using API to get summarized structure for chat {chat_id}")
                import requests
                endpoint = f"acra/{chat_id}/summarize_structure"
                
                # Prepare JSON payload for the API
                api_payload = {"timestamp": timestamp}
                if additional_info: 
                    api_payload["add_info"] = additional_info
                
                url = f"{acra_config.get('API_URL')}/{endpoint}"
                response = requests.post(url, json=api_payload)
                
                if response.status_code == 200:
                    summarized_json_data = response.json()
                    if "error" in summarized_json_data and summarized_json_data["error"]:
                        log.error(f"API summarization error for chat {chat_id}: {summarized_json_data['error']}")
                        return f"Erreur de l'API lors de la récupération de la structure résumée: {summarized_json_data['error']}"
                    
                    # Validate that the API returned a properly structured JSON response
                    if "projects" not in summarized_json_data:
                        log.error(f"API response for chat {chat_id} missing 'projects' key. Response: {summarized_json_data}")
                        return f"Réponse invalide de l'API: la clé 'projects' est manquante."

                    # The API returned a valid project structure, now generate a PowerPoint from it
                    temp_pptx_path_or_error = self._generate_summary_powerpoint(summarized_json_data, timestamp)
                    
                    if "error:" in temp_pptx_path_or_error.lower():
                        log.error(f"Failed to generate summary PPTX from API data for chat {chat_id}: {temp_pptx_path_or_error}")
                        return f"Erreur lors de la création du fichier PowerPoint de résumé: {temp_pptx_path_or_error.split('error:', 1)[-1].strip()}"
                    generated_pptx_path = temp_pptx_path_or_error
                else:
                    log.error(f"API call for summarized structure failed for chat {chat_id}: {response.status_code} - {response.text}")
                    return f"Erreur API ({response.status_code}) lors de la récupération de la structure résumée."
            
            else:
                # === DIRECT FUNCTION CALL WORKFLOW ===
                # In this path, core.summarize_ppt handles both summarization and PowerPoint generation
                log.info(f"Using direct call to summarize and generate PPTX for chat {chat_id}")
                raw_input_structure_for_llm: Optional[Dict[str, Any]] = None
                
                # Use cached structure if available for more efficient processing
                if self.cached_structure and isinstance(self.cached_structure, dict) and self.cached_structure.get("projects") is not None:
                    log.info(f"Using cached_structure for summarization for chat {chat_id}")
                    raw_input_structure_for_llm = self.cached_structure
                else:
                    log.info(f"No valid cached_structure found for chat {chat_id}. Summarize will process files.")
                
                # Call core.summarize_ppt, which handles:
                # 1. Processing files to extract data (if raw_structure_data is None)
                # 2. Summarizing the data with an LLM
                # 3. Generating a PowerPoint file from the summarized data
                # Returns either {"filename": "...", "summary": "/path/to/file.pptx"} or {"error": "message", "summary": None}
                result_from_core_summarize_ppt = summarize_ppt(
                    chat_id=chat_id, 
                    add_info=additional_info, 
                    timestamp=timestamp,
                    raw_structure_data=raw_input_structure_for_llm
                )
                
                # Handle possible error from summarize_ppt
                if "error" in result_from_core_summarize_ppt and result_from_core_summarize_ppt["error"]:
                    error_message = result_from_core_summarize_ppt["error"]
                    log.error(f"Error from summarize_ppt for chat {chat_id}: {error_message}")
                    return f"Erreur lors de la génération du résumé: {error_message}"
                
                # summarize_ppt has already generated the PPTX file, so we just need the path
                if "summary" not in result_from_core_summarize_ppt:
                    log.error(f"summarize_ppt didn't return a 'summary' key with file path: {result_from_core_summarize_ppt}")
                    return "Erreur: Le service de résumé n'a pas fourni le chemin du fichier généré."
                
                # The path to the already generated PPTX is in the "summary" key
                generated_pptx_path = result_from_core_summarize_ppt["summary"]

            # === COMMON WORKFLOW FOR BOTH PATHS ===
            # At this point, we should have a valid PowerPoint file path in generated_pptx_path
            if not generated_pptx_path:
                log.critical(f"Unexpectedly reached common logic with no generated_pptx_path for chat {chat_id}. This indicates a flaw in prior error trapping.")
                return "Erreur critique: Le chemin du fichier PowerPoint n'a pas été obtenu et l'erreur n'a pas été interceptée plus tôt."

            # Upload the generated PowerPoint file to OpenWebUI for user access
            upload_result = self.file_manager.upload_to_openwebui(generated_pptx_path)
            
            if "error" in upload_result:
                log.error(f"Résumé PPTX généré ({generated_pptx_path}) mais erreur d'upload pour chat {chat_id}: {upload_result['error']}")
                generated_filename = os.path.basename(generated_pptx_path) if generated_pptx_path else "inconnu"
                return f"Résumé généré ({generated_filename}) mais erreur lors du téléchargement vers OpenWebUI: {upload_result['error']}"
            
            # Build the success response message
            response_message_parts = []
            if self.system_prompt:
                try:
                    # Generate an introduction based on the system prompt if available
                    introduction = model_manager.generate_introduction(self.system_prompt)
                    response_message_parts.append(introduction)
                except Exception as intro_e:
                    log.warning(f"Could not generate introduction for chat {chat_id}: {intro_e}")
            
            generated_filename_for_msg = os.path.basename(generated_pptx_path) if generated_pptx_path else "résumé"
            response_message_parts.append(f"Le résumé ({generated_filename_for_msg}) a été généré avec succès.")
            response_message_parts.append(f"### URL de téléchargement:\n{upload_result.get('download_url', 'Non disponible')}")
            
            final_response = "\n\n".join(response_message_parts)
            
            # Save the file mapping to maintain state across sessions
            self.file_manager.save_file_mappings()
            return final_response
            
        except Exception as e:
            current_chat_id = self.file_manager.chat_id if hasattr(self, 'file_manager') and self.file_manager and self.file_manager.chat_id else "UNKNOWN_CHAT_ID"
            log.error(f"Exception in _execute_summarize for chat {current_chat_id}: {str(e)}", exc_info=True)
            return f"Erreur majeure lors de l'exécution du résumé: {str(e)}"
    
    def handle_structure_command(self) -> str:
        """
        Handle the /structure command, which analyzes and displays the structure of uploaded files.
        
        This method:
        1. Uses cached structure if available
        2. Otherwise, calls get_slide_structure to analyze PowerPoint files
        3. Formats the structure data into a human-readable format
        
        Returns:
            str: Formatted structure data or error message
        """
        try:
            if self.cached_structure is None:
                # Import here to avoid circular imports
                from core import get_slide_structure
                
                response = get_slide_structure(self.file_manager.chat_id)
                
                if "error" in response:
                    return f"Erreur lors de l'analyse de la structure: {response['error']}"
                
                # Cache the structure
                self.cached_structure = response
                # Format for display
                formatted_response = self._format_slide_data(response)
                return formatted_response
            else:
                # Use cached structure
                if isinstance(self.cached_structure, dict):
                    return self._format_slide_data(self.cached_structure)
                else:
                    return self.cached_structure
                    
        except Exception as e:
            log.error(f"Error handling structure command: {str(e)}")
            return f"Erreur lors de l'analyse de la structure: {str(e)}"
    
    def handle_generate_command(self, message: str) -> str:
        """Handle /generate command"""
        try:
            # Extract text content after the command
            text_content = message.replace("/generate", "").strip()
            if not text_content:
                return "Veuillez fournir du texte après la commande /generate pour générer un rapport."
            
            # Import here to avoid circular imports
            from core import generate_pptx_from_text
            
            # Generate timestamp for unique filename
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            
            if acra_config.get("USE_API"):
                # Use API endpoint
                import requests
                endpoint = f"generate_report/{self.file_manager.chat_id}?info={text_content}&timestamp={timestamp}"
                url = f"{acra_config.get('API_URL')}/{endpoint}"
                response = requests.get(url)
                result = response.json() if response.status_code == 200 else {"error": "Request failed"}
            else:
                # Use direct function call
                result = generate_pptx_from_text(self.file_manager.chat_id, text_content, timestamp)
            
            if "error" in result:
                return f"Erreur lors de la génération du rapport: {result['error']}"
            
            # Upload result and get download URL
            upload_result = self.file_manager.upload_to_openwebui(result["summary"])
            
            if "error" in upload_result:
                return f"Rapport généré mais erreur lors du téléchargement: {upload_result['error']}"
            
            self.file_manager.save_file_mappings()
            return f"Le rapport a été généré avec succès à partir du texte fourni.\n\n### URL de téléchargement:\n{upload_result.get('download_url', 'Non disponible')}"
            
        except Exception as e:
            log.error(f"Error handling generate command: {str(e)}")
            return f"Erreur lors de la génération du rapport: {str(e)}"
    
    def handle_clear_command(self, message: str) -> str:
        """Handle /clear command"""
        try:
            # Preserve current chat ID
            preserve_ids = [self.file_manager.chat_id] if self.file_manager.chat_id else []
            
            # Extract additional IDs to preserve if specified
            if " " in message:
                additional_ids = message.split(" ", 1)[1].strip().split()
                if additional_ids:
                    preserve_ids.extend(additional_ids)
            
            if acra_config.get("USE_API"):
                # Use API endpoint
                import requests
                endpoint = "acra/cleanup"
                url = f"{acra_config.get('API_URL')}/{endpoint}"
                
                # Send preserved IDs to the API
                payload = {"preserve_ids": preserve_ids}
                response = requests.post(url, json=payload)
                
                if response.status_code != 200:
                    return f"Erreur lors du nettoyage: API request failed with status {response.status_code}"
                
                cleanup_result = response.json()
            else:
                # Use direct function call
                cleanup_result = cleanup_orphaned_folders(preserve_ids=preserve_ids)
            
            # Reset state
            self.reset_state()
            self.file_manager.file_id_mapping = {}
            
            return f"Nettoyage terminé!\n\nProtégés: {preserve_ids}\nRésultat: {cleanup_result.get('message', 'Nettoyage effectué')}"
            
        except Exception as e:
            log.error(f"Error handling clear command: {str(e)}")
            return f"Une erreur s'est produite lors du nettoyage: {str(e)}"
    
    def handle_merge_command(self) -> str:
        """Handle /merge command"""
        try:
            chat_id = self.file_manager.chat_id
            if not chat_id:
                return "Error: No chat ID is set. Cannot merge files."

            if acra_config.get("USE_API"):
                # Use API endpoint
                import requests
                endpoint = f"acra/merge/{chat_id}"
                url = f"{acra_config.get('API_URL')}/{endpoint}"
                response = requests.post(url)
                merge_result = response.json() if response.status_code == 200 else {"error": f"API request failed with status {response.status_code}"}
            else:
                # Use direct function call
                # Import here to avoid circular imports
                from services import merge_pptx
                
                output_merge = os.path.join(acra_config.output_folder, chat_id, "merged")
                input_merge = acra_config.get_conversation_upload_folder(chat_id)
                
                merge_result = merge_pptx(input_merge, output_merge)
            
            if "error" in merge_result:
                return f"Erreur lors de la fusion des fichiers: {merge_result['error']}"
            
            # Get the merged file and upload to OpenWebUI
            merged_file = merge_result.get("merged_file")
            if merged_file:
                upload_result = self.file_manager.upload_to_openwebui(merged_file)
                if "error" in upload_result:
                    return f"Les fichiers ont été fusionnés avec succès, mais une erreur s'est produite lors de la génération du lien de téléchargement: {upload_result['error']}"
                else:
                    self.file_manager.save_file_mappings()
                    return f"Les fichiers ont été fusionnés avec succès.\n\n### URL de téléchargement:\n{upload_result.get('download_url', 'Non disponible')}"
            else:
                return "Les fichiers ont été fusionnés avec succès, mais le chemin du fichier fusionné n'a pas été trouvé."
                
        except Exception as e:
            log.error(f"Error handling merge command: {str(e)}")
            return f"Erreur lors de la fusion des fichiers: {str(e)}"
    
    def handle_regroup_command(self) -> str:
        """Handle /regroup command"""
        try:
            chat_id = self.file_manager.chat_id
            if not chat_id:
                return "Error: No chat ID is set. Cannot regroup projects."

            # Prepare payload and structure data
            cached_structure = None
            if self.cached_structure is not None and isinstance(self.cached_structure, dict) and "projects" in self.cached_structure:
                cached_structure = self.cached_structure

            if acra_config.get("USE_API"):
                # Use API endpoint
                import requests
                import json
                
                endpoint = f"acra/regroup/{chat_id}"
                url = f"{acra_config.get('API_URL')}/{endpoint}"
                
                # Prepare the payload with cached structure if available
                payload = {}
                if cached_structure:
                    payload["structure_data"] = cached_structure
                
                # Call the API
                response = requests.post(url, json=payload)
                if response.status_code != 200:
                    return f"Erreur lors de la réorganisation des données: API request failed with status {response.status_code}"
                
                regroup_result = response.json()
                
                # Update cached structure if the API returns a new structure
                if "structure" in regroup_result:
                    self.cached_structure = regroup_result["structure"]
                
                # Handle file path from API response
                result_path = regroup_result.get("path")
                if not result_path:
                    return "Les informations des projets ont été regroupées avec succès, mais le chemin du fichier n'a pas été trouvé."
                
                # Upload to OpenWebUI
                upload_result = self.file_manager.upload_to_openwebui(result_path)
            else:
                # Use direct function approach - get structure data first
                if cached_structure is None:
                    from core import get_slide_structure
                    structure_result = get_slide_structure(chat_id)
                    if "error" in structure_result:
                        return f"Erreur lors de l'analyse de la structure: {structure_result['error']}"
                else:
                    if isinstance(cached_structure, str):
                        from core import get_slide_structure
                        structure_result = get_slide_structure(chat_id)
                        if "error" in structure_result:
                            return f"Erreur lors de l'analyse de la structure: {structure_result['error']}"
                    else:
                        structure_result = cached_structure
                
                if not isinstance(structure_result, dict) or "projects" not in structure_result:
                    return f"Erreur: structure de données invalide. Type: {type(structure_result)}"
                
                # Get project grouping suggestions from LLM
                project_names = list(structure_result["projects"].keys())
                grouping_response = model_manager.generate_project_grouping(project_names)
                
                try:
                    groups_to_merge = extract_json(grouping_response)
                    if not isinstance(groups_to_merge, list):
                        groups_to_merge = []
                except:
                    log.warning("Could not extract valid JSON from LLM response")
                    groups_to_merge = []
                
                log.info(f"Groups to merge: {groups_to_merge}")
                
                # Process regrouping
                new_structure = self._process_regrouping(structure_result, groups_to_merge)
                
                # Generate PowerPoint with regrouped data
                result = self._generate_regrouped_powerpoint(new_structure)
                
                # Update cached structure
                self.cached_structure = new_structure
                
                return result
            
            # Common upload handling code for API path
            if "error" in upload_result:
                return f"Les informations des projets ont été regroupées avec succès, mais une erreur s'est produite lors de la génération du lien de téléchargement: {upload_result['error']}"
            
            self.file_manager.save_file_mappings()
            return f"Les informations des projets ont été regroupées avec succès.\n\n### URL de téléchargement:\n{upload_result.get('download_url', 'Non disponible')}"
            
        except Exception as e:
            log.error(f"Error handling regroup command: {str(e)}")
            return f"Erreur lors de la réorganisation des données: {str(e)}"
    
    def _format_slide_data(self, data: dict) -> str:
        """
        Format slide structure data into a readable text format.
        
        This method processes the raw structure data from PowerPoint files into a
        hierarchical, Markdown-formatted text display. It includes:
        - Project hierarchies with proper indentation
        - Icons to distinguish different levels of projects/subprojects
        - Critical alerts, minor alerts, and advancements
        - Upcoming events organized by service
        
        Args:
            data (dict): Raw structure data from PowerPoint analysis
            
        Returns:
            str: Markdown-formatted text representation of the structure
        """
        if not data:
            return "Aucun fichier PPTX fourni."
        
        projects = data.get("projects", {})
        if not projects:
            return "Aucun projet trouvé dans les fichiers analysés."
        
        metadata = data.get("metadata", {})
        processed_files = metadata.get("processed_files", 0)
        upcoming_events = data.get("upcoming_events", {})
        
        def format_project_hierarchy(project_name, content, level=0):
            """
            Recursively format a project and its subprojects with proper indentation and styling.
            
            Args:
                project_name (str): Name of the project
                content (dict): Project content data
                level (int): Indentation level (0 for top-level)
                
            Returns:
                str: Formatted project text
            """
            output = ""
            indent = "  " * level
            
            # Format project name based on level
            if level == 0:
                output += f"{indent}🔶 **{project_name}**\n"
            elif level == 1:
                output += f"{indent}📌 **{project_name}**\n"
            else:
                output += f"{indent}📎 *{project_name}*\n"
            
            # Add project information
            if "information" in content and content["information"]:
                info_lines = content["information"].split('\n')
                for line in info_lines:
                    if line.strip():
                        output += f"{indent}- {line}\n"
                output += "\n"
            
            # Add critical alerts
            if "critical" in content and content["critical"]:
                output += f"{indent}- 🔴 **Alertes Critiques:**\n"
                for alert in content["critical"]:
                    output += f"{indent}  - {alert}\n"
                output += "\n"
            
            # Add minor alerts
            if "small" in content and content["small"]:
                output += f"{indent}- 🟡 **Alertes à surveiller:**\n"
                for alert in content["small"]:
                    output += f"{indent}  - {alert}\n"
                output += "\n"
            
            # Add advancements
            if "advancements" in content and content["advancements"]:
                output += f"{indent}- 🟢 **Avancements:**\n"
                for advancement in content["advancements"]:
                    output += f"{indent}  - {advancement}\n"
                output += "\n"
            
            # Recursively process subprojects
            for key, value in content.items():
                if isinstance(value, dict) and key not in ["information", "critical", "small", "advancements"]:
                    output += format_project_hierarchy(key, value, level + 1)
            
            return output
        
        result = f"📊 **Synthèse globale de {processed_files} fichier(s) analysé(s)**\n\n"
        
        # Add projects
        for project_name, project_content in projects.items():
            result += format_project_hierarchy(project_name, project_content)
        
        # Add upcoming events
        if upcoming_events:
            result += "\n\n📅 **Événements à venir par service:**\n\n"
            for service, events in upcoming_events.items():
                if events:
                    result += f"- **{service}:**\n"
                    for event in events:
                        result += f"  - {event}\n"
                    result += "\n"
        else:
            result += "\n\n📅 **Événements à venir:** Aucun événement particulier prévu.\n"
        
        return result.strip()
    
    def _process_regrouping(self, structure_result: dict, groups_to_merge: list) -> dict:
        """Process the regrouping of projects"""
        new_structure = json.loads(json.dumps(structure_result))
        
        for group in groups_to_merge:
            if not isinstance(group, list) or len(group) < 2:
                continue
            
            # Create mapping between original and cleaned names
            original_to_cleaned = {}
            cleaned_to_original = {}
            
            for original_name in group:
                cleaned_name = original_name.replace('\n', ' ').strip()
                original_to_cleaned[original_name] = cleaned_name
                cleaned_to_original[cleaned_name] = original_name
            
            cleaned_group = list(cleaned_to_original.keys())
            main_project_cleaned = min(cleaned_group, key=len)
            main_project_original = cleaned_to_original[main_project_cleaned]
            
            other_projects_cleaned = [p for p in cleaned_group if p != main_project_cleaned]
            other_projects_original = [cleaned_to_original[p] for p in other_projects_cleaned 
                                     if cleaned_to_original[p] in new_structure["projects"]]
            
            if main_project_original not in new_structure["projects"]:
                log.warning(f"Main project '{main_project_original}' not found in structure, skipping group")
                continue
            
            # Regroup other projects under main project
            for other_project_original in other_projects_original:
                if other_project_original in new_structure["projects"]:
                    try:
                        other_data = new_structure["projects"][other_project_original]
                        other_project_cleaned = original_to_cleaned[other_project_original]
                        
                        sub_name = other_project_cleaned.replace(main_project_cleaned, "").strip("_").strip()
                        if not sub_name:
                            sub_name = other_project_cleaned
                        
                        # Handle terminal vs non-terminal projects
                        if "information" in new_structure["projects"][main_project_original]:
                            main_data = {
                                "information": new_structure["projects"][main_project_original].get("information", ""),
                                "critical": new_structure["projects"][main_project_original].get("critical", []),
                                "small": new_structure["projects"][main_project_original].get("small", []),
                                "advancements": new_structure["projects"][main_project_original].get("advancements", [])
                            }
                            
                            new_structure["projects"][main_project_original] = {
                                "Général": main_data,
                                sub_name: other_data
                            }
                        else:
                            new_structure["projects"][main_project_original][sub_name] = other_data
                        
                        del new_structure["projects"][other_project_original]
                        log.info(f"Moved {other_project_original} to {main_project_original}.{sub_name}")
                        
                    except Exception as e:
                        log.error(f"Error moving project {other_project_original}: {str(e)}")
                        continue
        
        return new_structure
    
    def _generate_regrouped_powerpoint(self, new_structure: dict) -> str:
        """Generate PowerPoint with regrouped data"""
        try:
            from src.services.update_pttx_service import update_table_with_project_data
            from pptx import Presentation
            from pptx.util import Pt
            
            # Create output directory
            output_regroup = os.path.join(acra_config.output_folder, self.file_manager.chat_id, "regrouped")
            os.makedirs(output_regroup, exist_ok=True)
            
            # Generate output filename with timestamp
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            output_file = os.path.join(output_regroup, f"regrouped_{timestamp}.pptx")
            
            # Create presentation from template or blank
            if os.path.exists(acra_config.template_path):
                prs = Presentation(acra_config.template_path)
            else:
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                table_shape = slide.shapes.add_table(rows=10, cols=3, left=Pt(30), top=Pt(30), width=Pt(600), height=Pt(400))
            
            # Save temporary file
            temp_path = os.path.join(output_regroup, "temp.pptx")
            prs.save(temp_path)
            
            # Update with project data
            updated_path = update_table_with_project_data(
                temp_path,
                0,  # slide index
                0,  # table shape index
                new_structure["projects"],
                output_file,
                new_structure.get("upcoming_events", {})
            )
            
            # Clean up temp file
            if os.path.exists(temp_path):
                os.remove(temp_path)
            
            # Upload to OpenWebUI
            upload_result = self.file_manager.upload_to_openwebui(updated_path)
            
            if "error" in upload_result:
                return f"Les informations des projets ont été regroupées avec succès, mais une erreur s'est produite lors de la génération du lien de téléchargement: {upload_result['error']}"
            
            self.file_manager.save_file_mappings()
            return f"Les informations des projets ont été regroupées avec succès.\n\n### URL de téléchargement:\n{upload_result.get('download_url', 'Non disponible')}"
            
        except Exception as e:
            log.error(f"Error generating regrouped PowerPoint: {str(e)}")
            return f"Erreur lors de la génération de la présentation PowerPoint: {str(e)}" 