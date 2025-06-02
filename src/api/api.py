import os,sys
import uvicorn
import logging
import datetime
from fastapi import FastAPI, HTTPException
from fastapi.responses import FileResponse  
from dotenv import load_dotenv
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))
from core import summarize_ppt, get_slide_structure, get_slide_structure_wcolor, delete_all_pptx_files, generate_pptx_from_text
from services import merge_pptx
from services.cleanup_service import cleanup_orphaned_folders

app = FastAPI() 
load_dotenv()
logger = logging.getLogger(__name__)
UPLOAD_FOLDER = os.getenv("UPLOAD_FOLDER", "pptx_folder")
OUTPUT_FOLDER = os.getenv("OUTPUT_FOLDER", "OUTPUT")

@app.get("/acra/{folder_name}?add_info={add_info}")
async def summarize(folder_name: str, add_info: str = None):
    """
    Summarizes the content of PowerPoint files in a folder and updates a template PowerPoint file with the summary.
    
    This endpoint:
    1. Takes a folder name containing PowerPoint files
    2. Processes all PowerPoint files in that folder
    3. Extracts and summarizes key information using an LLM
    4. Generates a new PowerPoint with the summarized content
    
    The PowerPoint will be structured with a hierarchical format:
      - Main projects as headers
      - Subprojects under each main project
      - Information, alerts for each subproject
      - Events listed by service at the bottom

    Args:
        folder_name (str): The name of the folder containing PowerPoint files to analyze.
        add_info (str, optional): Additional information or instructions for guiding the summarization process.

    Returns:
        dict: A dictionary containing the filename and path of the generated PowerPoint file.
              Format: {"filename": str, "summary": str}

    Raises:
        HTTPException: If there's an error processing the PowerPoint files.
    """
    logger.info(f"Summarizing PPT for folder: {folder_name}")
    try:
        return summarize_ppt(folder_name, add_info)
    except Exception as e:
        # Log the exception for debugging
        print(f"Error in summarize_folder: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Summarize error: {str(e)}")

@app.post("/acra/generate_report/{folder_name}?info={info}")
async def generate_report(folder_name: str, info: str):
    """
    Generates a PowerPoint report from text input using GET request.
    
    This endpoint:
    1. Takes text information and a folder name
    2. Uses an LLM to extract structured data from the text
    3. Creates a PowerPoint presentation based on the extracted structure
    
    Args:
        folder_name (str): The name of the folder to store the generated PowerPoint
        info (str): Text information to analyze and convert to PowerPoint
        
    Returns:
        dict: A dictionary containing the filename and path of the generated PowerPoint
              Format: {"filename": str, "summary": str}
        
    Raises:
        HTTPException: If there's an error generating the PowerPoint
    """
    logger.info(f"Generating report from text (GET) for folder: {folder_name}")
    try:
        return generate_pptx_from_text(folder_name, info)
    except Exception as e:
        print(f"Error in generate_report (GET): {str(e)}")
        raise HTTPException(status_code=500, detail=f"Report generation error: {str(e)}")
    
@app.get("/get_slide_structure/{foldername}")
async def get_structure(foldername: str):
    """
    Analyzes all PowerPoint files in a folder and extracts their structured content.
    
    This endpoint:
    1. Finds all PowerPoint files in the specified folder
    2. Extracts project data and upcoming events from each file
    3. Merges the data from all files into a coherent structure
    
    The resulting structure is used for summarization, analysis, and PowerPoint generation.
    It forms the foundation of the ACRA system's understanding of presentation content.
    
    Args:
        foldername (str): Name of the folder containing PowerPoint files to analyze
        
    Returns:
        dict: A structured representation of all PowerPoint content with:
              - projects: Hierarchical project data
              - upcoming_events: Events organized by service
              - metadata: Processing information
              - source_files: Details about processed files
        
    Raises:
        HTTPException: If there's an error analyzing the presentations
    """
    logger.info(f"Getting slide structure for folder: {foldername}")
    try:
        return get_slide_structure(foldername)
    except Exception as e:
        print(f"Error in slide_structure: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Slide structure error: {str(e)}")

# Endpoint for getting slide structure with color information
@app.get("/get_slide_structure_wcolor/{filename}")
async def structure_wcolor(filename: str):
    """
    Analyzes a single PowerPoint file with color extraction.
    
    This endpoint uses a different analysis method that preserves color information
    from the slides, which can be useful for certain visualization or analysis tasks.
    Unlike the regular structure endpoint, this operates on a single file rather than
    a folder of files.
    
    Args:
        filename (str): The name of the PowerPoint file to analyze
        
    Returns:
        dict: Analysis results with color information, including:
              - filename: Name of the analyzed file
              - slide_data: Extracted slide content with color information
        
    Raises:
        HTTPException: If the file is not found or cannot be processed
    """
    logger.info(f"Getting slide structure with color for file: {filename}")
    try:
        return get_slide_structure_wcolor(filename)
    except Exception as e:
        print(f"Error in slide_structure_wcolor: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Slide structure with color error: {str(e)}")

# Endpoint for merging PowerPoint files
@app.post("/acra/merge/{chat_id}")
async def merge_files(chat_id: str):
    """
    Merges all PowerPoint files in a folder into a single presentation.
    
    This endpoint:
    1. Identifies all PowerPoint files in the specified chat/folder
    2. Combines them into a single PowerPoint presentation
    3. Preserves the structure and content of each original file
    
    This is useful for creating a comprehensive presentation from multiple
    individual presentations, such as combining reports from different teams.
    
    Args:
        chat_id (str): Identifier for the chat/folder containing files to merge
        
    Returns:
        dict: A dictionary with information about the merged file, including:
              - merged_file: Path to the merged PowerPoint file
              - message: Status message
        
    Raises:
        HTTPException: If there's an error during the merge operation
    """
    logger.info(f"Merging PowerPoint files for chat: {chat_id}")
    try:
        # Set up input and output directories
        input_dir = os.path.join(UPLOAD_FOLDER, chat_id)
        output_dir = os.path.join(OUTPUT_FOLDER, chat_id, "merged")
        os.makedirs(output_dir, exist_ok=True)
        
        # Merge the files
        result = merge_pptx(input_dir, output_dir)
        
        if "error" in result:
            raise HTTPException(status_code=500, detail=result["error"])
            
        return result
    except HTTPException as he:
        raise he
    except Exception as e:
        print(f"Error merging files: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Merge error: {str(e)}")

# Endpoint for regrouping project information
@app.post("/acra/regroup/{chat_id}")
async def regroup_projects(chat_id: str, body: dict = None):
    """
    Regroups project information by combining similar or related projects.
    
    This endpoint:
    1. Analyzes the structure of PowerPoint files in the specified chat/folder
    2. Identifies projects with similar themes or names
    3. Suggests groupings for these projects
    4. Creates a new PowerPoint with the regrouped information
    
    This is particularly useful for organizing information when multiple
    presentations have overlapping topics or when a topic is spread across
    different presentations.
    
    Args:
        chat_id (str): Identifier for the chat/folder containing files to analyze
        body (dict, optional): Optional configuration parameters:
            - structure_data (dict, optional): Pre-extracted structure data
            - groups_to_merge (list, optional): Explicit groups to merge
        
    Returns:
        dict: A dictionary with information about the regrouped data, including:
              - filename: Name of the generated PowerPoint file
              - path: Path to the generated PowerPoint file
        
    Raises:
        HTTPException: If there's an error during the regrouping process
    """
    logger.info(f"Regrouping projects for chat: {chat_id}")
    try:
        # First get the structure data
        structure_data = None
        if body and "structure_data" in body:
            structure_data = body["structure_data"]
        else:
            # Get structure from files
            structure_data = get_slide_structure(chat_id)
        
        if "error" in structure_data or not isinstance(structure_data, dict) or "projects" not in structure_data:
            error_msg = "Invalid structure data for regrouping"
            if isinstance(structure_data, dict) and "error" in structure_data:
                error_msg = structure_data["error"]
            raise HTTPException(status_code=400, detail=error_msg)
        
        # Extract project names for grouping
        from OLLibrary.utils.json_service import extract_json
        from services.model_manager import model_manager
        
        project_names = list(structure_data["projects"].keys())
        
        # Get grouping from request or generate via LLM
        groups_to_merge = []
        if body and "groups_to_merge" in body and isinstance(body["groups_to_merge"], list):
            groups_to_merge = body["groups_to_merge"]
        else:
            # Generate grouping suggestions
            grouping_response = model_manager.generate_project_grouping(project_names)
            try:
                groups_to_merge = extract_json(grouping_response)
                if not isinstance(groups_to_merge, list):
                    groups_to_merge = []
            except:
                logger.warning("Could not extract valid JSON from LLM response")
        
        # Process regrouping
        from src.services.command_handler import CommandHandler
        from src.services.file_manager import FileManager
        
        # Create temporary handlers to use their processing methods
        temp_file_manager = FileManager(chat_id=chat_id)
        temp_handler = CommandHandler(file_manager=temp_file_manager)
        
        # Use the handler's methods to process regrouping
        new_structure = temp_handler._process_regrouping(structure_data, groups_to_merge)
        
        # Generate PowerPoint with regrouped data
        output_dir = os.path.join(OUTPUT_FOLDER, chat_id, "regrouped")
        os.makedirs(output_dir, exist_ok=True)
        
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        output_file = os.path.join(output_dir, f"regrouped_{timestamp}.pptx")
        
        # Create the regrouped PowerPoint
        from src.services.update_pttx_service import update_table_with_project_data
        from pptx import Presentation
        from pptx.util import Pt
        
        # Create presentation from template or blank
        template_path = os.getenv("TEMPLATE_FILE", "templates/CRA_TEMPLATE_IA.pptx")
        if not os.path.isabs(template_path):
            template_path = os.path.join(os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..")), template_path)
            
        if os.path.exists(template_path):
            prs = Presentation(template_path)
        else:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            table_shape = slide.shapes.add_table(rows=10, cols=3, left=Pt(30), top=Pt(30), width=Pt(600), height=Pt(400))
        
        # Save temporary file
        temp_path = os.path.join(output_dir, f"temp_{timestamp}.pptx")
        prs.save(temp_path)
        
        # Update with project data
        updated_path = update_table_with_project_data(
            temp_path,
            0,  # slide index
            0,  # table shape index
            new_structure["projects"],
            output_file,
            new_structure.get("upcoming_events", {})
        )
        
        # Clean up temp file
        if os.path.exists(temp_path):
            os.remove(temp_path)
            
        if "error" in updated_path.lower():
            raise HTTPException(status_code=500, detail=f"Failed to generate regrouped PowerPoint: {updated_path}")
            
        return {
            "filename": os.path.basename(updated_path),
            "path": updated_path,
            "message": "Projects successfully regrouped"
        }
        
    except HTTPException as he:
        raise he
    except Exception as e:
        print(f"Error regrouping projects: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Regroup error: {str(e)}")

# Endpoint for downloading files
@app.get("/download/{folder_name}/{filename}")
async def download_file(folder_name: str, filename: str):
    """
    Downloads a file from a specific folder on the server.
    
    This endpoint allows clients to download generated PowerPoint files or other
    output files. It's typically used after operations that create files,
    such as summarization or text-to-PowerPoint generation.
    
    Args:
        folder_name (str): The name of the folder containing the file
        filename (str): The name of the file to be downloaded
        
    Returns:
        FileResponse: A response containing the file to be downloaded,
                      with appropriate content type headers
        
    Raises:
        HTTPException: If the file does not exist (404 status code)
    """
    logger.info(f"Download request for file: {filename} in folder: {folder_name}")
    file_path = os.path.join(os.getenv("OUTPUT_FOLDER", "OUTPUT"), folder_name, filename)
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail=f"File not found at path: {file_path}")
    
    return FileResponse(
        path=file_path,
        filename=filename,
        media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )

# Endpoint for deleting all PowerPoint files in a folder
@app.delete("/delete_all_pptx_files/{foldername}")
async def delete_files(foldername: str):
    """
    Deletes all PowerPoint files in the specified folder.
    
    This endpoint is used for cleanup operations when files are no longer needed.
    It removes all .pptx files from the specified folder in the upload directory.
    
    Args:
        foldername (str): The name of the folder containing files to delete
        
    Returns:
        dict: Message indicating the result of the operation
              Format: {"message": str}
        
    Raises:
        HTTPException: If the folder doesn't exist or files can't be deleted
    """
    logger.info(f"Deleting all PPTX files in folder: {foldername}")
    try:
        return delete_all_pptx_files(foldername)
    except Exception as e:
        print(f"Error during file deletion: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Deletion error: {str(e)}")

# Endpoint for cleaning up orphaned conversations
@app.post("/acra/cleanup")
async def cleanup_orphaned(body: dict = None):
    """
    Cleans up orphaned conversation folders and files.
    
    This endpoint:
    1. Identifies conversation folders that are no longer in use
    2. Deletes these folders and their associated files
    3. Preserves specified conversations
    
    This is useful for system maintenance to free up disk space and remove
    files that are no longer needed. It's typically called periodically or
    when switching between conversations.
    
    Args:
        body (dict, optional): Optional JSON body containing:
            - preserve_ids (list): List of conversation IDs to preserve
            
    Returns:
        dict: Results of the cleanup operation, including:
              - deleted_folders: Number of folders deleted
              - deleted_files: Number of files deleted
              - message: Summary of the operation
        
    Raises:
        HTTPException: If there's an error during the cleanup process
    """
    logger.info("Cleaning up orphaned conversations")
    try:
        preserve_ids = []
        if body and "preserve_ids" in body and isinstance(body["preserve_ids"], list):
            preserve_ids = body["preserve_ids"]
            logger.info(f"Will preserve the following conversation IDs: {preserve_ids}")
        
        # Call the cleanup function
        cleanup_result = cleanup_orphaned_folders(preserve_ids=preserve_ids)
        
        if not cleanup_result or not isinstance(cleanup_result, dict):
            raise HTTPException(status_code=500, detail="Cleanup function did not return a valid result")
            
        return cleanup_result
    except Exception as e:
        print(f"Error during orphaned conversation cleanup: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Cleanup error: {str(e)}")

# Function to run the API server
def run():
    """
    Runs the FastAPI server using uvicorn.
    
    This function is called when the module is run directly or imported
    and the run() function is explicitly called.
    """
    uvicorn.run(app, host="0.0.0.0", port=5050)


if __name__ == "__main__":
    # Example usage when module is run directly
    summarize_ppt("669fa53b-649a-4023-8066-4cd86670e88b")